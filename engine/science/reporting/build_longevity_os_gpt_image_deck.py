#!/usr/bin/env python3
"""Build the Longevity OS 3-slide intro deck from GPT-generated images."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "decks" / "longevity_os_intro"
IMG_DIR = OUT_DIR / "gpt_images"
OUT = OUT_DIR / "longevity_os_intro_gpt_images_3slides.pptx"
NOTES = OUT_DIR / "longevity_os_intro_gpt_images_speaker_notes.md"


WHITE = RGBColor(255, 255, 255)
INK = RGBColor(15, 23, 42)
MUTED = RGBColor(82, 92, 110)
BLUE = RGBColor(0, 68, 120)
GREEN = RGBColor(0, 121, 89)
ORANGE = RGBColor(198, 92, 0)
LIGHT = RGBColor(231, 236, 243)


def add_text(slide, text, x, y, w, h, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_rule(slide, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.035)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_header(slide, title, subtitle, number, accent):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    add_text(slide, "Longevity OS", 0.32, 0.13, 1.5, 0.18, 7.5, MUTED, True)
    add_text(slide, number, 12.4, 0.13, 0.45, 0.18, 7.5, MUTED, False, PP_ALIGN.RIGHT)
    add_text(slide, title, 0.32, 0.35, 10.0, 0.32, 17, INK, True)
    add_text(slide, subtitle, 0.32, 0.73, 10.8, 0.25, 8.8, MUTED)
    add_rule(slide, 0.32, 0.98, 1.25, accent)


def add_footer(slide, text):
    add_text(slide, text, 0.35, 7.16, 12.6, 0.18, 6.3, MUTED)


def add_image(slide, image_name):
    # GPT images are generated as 16:9. Place them as the main panel with
    # preserved aspect ratio beneath the header.
    slide.shapes.add_picture(
        str(IMG_DIR / image_name), Inches(0.42), Inches(1.18), width=Inches(12.5)
    )


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_header(
        slide,
        "Motivation: aging is the hidden state behind divergent health trajectories",
        "Chronological age is observable; reserve, vulnerability, and multisystem disease risk are latent.",
        "01",
        BLUE,
    )
    add_image(slide, "slide1_aging_state.png")
    add_footer(
        slide,
        "Framing: move beyond age as time elapsed toward aging as a measurable systems state.",
    )

    slide = prs.slides.add_slide(blank)
    add_header(
        slide,
        "Current status: biological age is measured by increasingly high-dimensional clocks",
        "Clinical, molecular, organ, and digital signals can be compressed into age or risk estimates.",
        "02",
        GREEN,
    )
    add_image(slide, "slide2_measurement_landscape.png")
    add_footer(
        slide,
        "Anchor examples: Horvath DNAm age, PhenoAge, GrimAge, DunedinPACE, organ-age and wearable-derived clocks.",
    )

    slide = prs.slides.add_slide(blank)
    add_header(
        slide,
        "Deficiency: clocks are often static snapshots of an interacting system",
        "The missing layer is dynamic coordination: lag, phase, feedback, perturbation response, and recovery.",
        "03",
        ORANGE,
    )
    add_image(slide, "slide3_dynamic_coordination.png")
    add_footer(
        slide,
        "Longevity OS thesis: measure aging as multimodal state plus cross-system coordination and resilience.",
    )

    prs.save(OUT)


def write_notes():
    NOTES.write_text(
        """# Longevity OS intro deck - GPT image version

## Slide 1

Open with the idea that age is a proxy. The relevant object is a latent physiological state: reserve, repair, inflammation, metabolic flexibility, vascular integrity, tissue structure, and stress response. Two people can be the same chronological age while occupying very different physiological trajectories.

## Slide 2

Give the field credit. Aging measurement has progressed from single biomarkers to composite clinical clocks, DNA methylation clocks, organ clocks, and digital physiology. These methods made aging measurable, but they often compress rich biology into a supervised age or risk estimate.

## Slide 3

State the deficiency directly: most clocks are snapshots. They capture burden and level, but they do not fully capture dynamics, context, interaction, lag, phase, perturbation response, and recovery. This sets up Longevity OS as a system for multimodal aging state plus coordination/resilience.

## Transition

The next slide should introduce AI-READI as the dataset that makes this possible: clinical labs, retinal imaging, ECG, CGM, wearable heart rate/activity/sleep, and environmental exposure measured in the same participants.
"""
    )


def main():
    missing = [
        p
        for p in [
            IMG_DIR / "slide1_aging_state.png",
            IMG_DIR / "slide2_measurement_landscape.png",
            IMG_DIR / "slide3_dynamic_coordination.png",
        ]
        if not p.exists()
    ]
    if missing:
        raise FileNotFoundError(f"Missing generated image assets: {missing}")
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build()
    write_notes()
    print(f"Wrote {OUT}")
    print(f"Wrote {NOTES}")


if __name__ == "__main__":
    main()
