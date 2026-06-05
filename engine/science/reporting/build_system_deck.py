"""
Build PPTX deck focused on the Hypothesis Discovery System design.
Visualizes: workflow, agents, lifecycle, feedback loops, hypothesis table,
theory vs data-driven, comparison with SOTA.
"""

from __future__ import annotations

import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as pe
import numpy as np
import pandas as pd
import scienceplots  # noqa
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Palette
C_BLUE = "#0072B2"
C_SKY  = "#56B4E9"
C_AMBER = "#E69F00"
C_VERM = "#D55E00"
C_GREEN = "#009E73"
C_PINK = "#CC79A7"
C_GRAY = "#999999"
C_DARK = "#1a1a2e"

FIGDIR = Path("results/figures/system_deck")
FIGDIR.mkdir(parents=True, exist_ok=True)
DPI = 300

plt.style.use(["science", "nature", "no-latex"])
plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Arial", "DejaVu Sans"],
    "font.size": 10,
    "axes.spines.top": False,
    "axes.spines.right": False,
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "savefig.facecolor": "white",
})


def save(fig, name):
    for fmt in ("png", "svg", "pdf"):
        fig.savefig(FIGDIR / f"{name}.{fmt}", format=fmt, dpi=DPI,
                    bbox_inches="tight", pad_inches=0.12)
    plt.close(fig)
    return str(FIGDIR / f"{name}.png")


def _box(ax, x, y, w, h, color, title, body="", title_size=11, body_size=8):
    rect = mpatches.FancyBboxPatch(
        (x - w/2, y - h/2), w, h,
        boxstyle="round,pad=0.08", facecolor="white",
        edgecolor=color, linewidth=2.0
    )
    ax.add_patch(rect)
    ax.text(x, y + h*0.2, title, ha="center", va="center",
            fontsize=title_size, fontweight="bold", color=color)
    if body:
        ax.text(x, y - h*0.15, body, ha="center", va="center",
                fontsize=body_size, color="#555555", linespacing=1.2)


def _arrow(ax, x1, y1, x2, y2, color="#aaaaaa", style="-|>", lw=1.5):
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle=style, color=color, lw=lw, mutation_scale=12))


# ══════════════════════════════════════════════════════════════════════
# Fig 1: Full pipeline workflow
# ══════════════════════════════════════════════════════════════════════

def fig_pipeline():
    fig, ax = plt.subplots(figsize=(12, 7))
    ax.set_xlim(-1, 13)
    ax.set_ylim(-2.5, 7)
    ax.axis("off")

    # Top: Dual-mode generation
    _box(ax, 1.5, 6, 2.5, 1.4, C_BLUE, "Theory-Driven",
         "Literature\nMechanisms\nKnown timescales")
    _box(ax, 4.5, 6, 2.5, 1.4, C_AMBER, "Data-Driven",
         "Null results\nSurprises\nAnomaly patterns")

    # Proposer
    _box(ax, 3, 3.8, 3, 1.4, C_BLUE, "PROPOSER AGENT",
         "Generates structured\nhypotheses (JSON)")
    ax.text(4.7, 3.8, "LLM", fontsize=7, color=C_BLUE,
            bbox=dict(boxstyle="round,pad=0.15", facecolor="#dbeafe", edgecolor=C_BLUE, lw=0.5))

    _arrow(ax, 1.5, 5.25, 2.2, 4.55, C_BLUE)
    _arrow(ax, 4.5, 5.25, 3.8, 4.55, C_AMBER)

    # Critic
    _box(ax, 7.5, 3.8, 3, 1.4, C_AMBER, "CRITIC AGENT",
         "10 constraint checks\nDimensional scoring")
    ax.text(9.2, 3.8, "LLM", fontsize=7, color=C_AMBER,
            bbox=dict(boxstyle="round,pad=0.15", facecolor="#fef3c7", edgecolor=C_AMBER, lw=0.5))

    _arrow(ax, 4.5, 3.8, 5.95, 3.8, C_GRAY)
    ax.text(5.2, 4.05, "propose", fontsize=8, color=C_GRAY, style="italic")

    # Reject path
    ax.annotate("", xy=(10.5, 5.5), xytext=(9.0, 4.55),
                arrowprops=dict(arrowstyle="-|>", color=C_VERM, lw=1.2, mutation_scale=10))
    ax.text(10.5, 5.7, "REJECTED", fontsize=9, color=C_VERM, fontweight="bold", ha="center")

    # Workspace (center)
    rect = mpatches.FancyBboxPatch((2, 1.0), 8, 1.6, boxstyle="round,pad=0.12",
                                    facecolor="#f8f8f8", edgecolor=C_GRAY, linewidth=1.5)
    ax.add_patch(rect)
    ax.text(6, 2.25, "HYPOTHESIS WORKSPACE", fontsize=11, fontweight="bold",
            ha="center", color=C_DARK)

    # Lifecycle pills
    states = [("proposed", C_SKY), ("critiqued", C_AMBER), ("validated", C_GREEN),
              ("executing", C_BLUE), ("completed", C_PINK), ("verified", C_DARK)]
    for i, (label, color) in enumerate(states):
        x = 2.8 + i * 1.35
        pill = mpatches.FancyBboxPatch((x - 0.55, 1.2), 1.1, 0.45,
                                        boxstyle="round,pad=0.08", facecolor=color, alpha=0.8)
        ax.add_patch(pill)
        ax.text(x, 1.42, label, fontsize=7, color="white", ha="center", va="center", fontweight="bold")
        if i < len(states) - 1:
            ax.annotate("", xy=(x + 0.6, 1.42), xytext=(x + 0.8, 1.42),
                        arrowprops=dict(arrowstyle="-|>", color="#cccccc", lw=0.8, mutation_scale=8))

    _arrow(ax, 7.5, 3.05, 7.5, 2.65, C_GREEN)
    ax.text(7.8, 2.85, "validate", fontsize=8, color=C_GREEN, style="italic")

    # Executor
    _box(ax, 3.5, -0.7, 3, 1.4, C_GREEN, "EXECUTOR AGENT",
         "Generate script\nSlurm 32-shard\nCollect results")
    ax.text(5.2, -0.7, "LLM+Code", fontsize=7, color=C_GREEN,
            bbox=dict(boxstyle="round,pad=0.15", facecolor="#d1fae5", edgecolor=C_GREEN, lw=0.5))

    _arrow(ax, 4.5, 0.95, 3.8, 0.05, C_GREEN)

    # Verifier
    _box(ax, 8.5, -0.7, 3, 1.4, C_PINK, "VERIFIER AGENT",
         "9 rule-based checks\n+ LLM interpretation")
    ax.text(10.2, -0.7, "Rules+LLM", fontsize=7, color=C_PINK,
            bbox=dict(boxstyle="round,pad=0.15", facecolor="#fce7f3", edgecolor=C_PINK, lw=0.5))

    _arrow(ax, 5.05, -0.7, 6.95, -0.7, C_GRAY)
    ax.text(6.0, -0.45, "completed", fontsize=8, color=C_GRAY, style="italic")

    # Verdict outputs
    verdicts = [("supported", C_GREEN, 7.5), ("refuted", C_VERM, 9.0), ("inconclusive", C_GRAY, 10.5)]
    for label, color, x in verdicts:
        ax.text(x, -2.0, label, fontsize=9, fontweight="bold", color=color, ha="center")
    _arrow(ax, 8.0, -1.45, 7.5, -1.75, C_GREEN, lw=1.0)
    _arrow(ax, 8.5, -1.45, 9.0, -1.75, C_VERM, lw=1.0)
    _arrow(ax, 9.0, -1.45, 10.5, -1.75, C_GRAY, lw=1.0)

    # Feedback loops
    ax.annotate("", xy=(1.0, 5.25), xytext=(7.5, -1.8),
                arrowprops=dict(arrowstyle="-|>", color=C_VERM, lw=1.2,
                                connectionstyle="arc3,rad=0.5", linestyle="--", mutation_scale=10))
    ax.text(-0.5, 1.5, "Null result\nlearning", fontsize=8, color=C_VERM, ha="center", style="italic")

    ax.annotate("", xy=(4.0, 5.25), xytext=(9.5, -1.8),
                arrowprops=dict(arrowstyle="-|>", color=C_GREEN, lw=1.2,
                                connectionstyle="arc3,rad=0.6", linestyle="--", mutation_scale=10))
    ax.text(12, 1.5, "Hypothesis\ndeepening", fontsize=8, color=C_GREEN, ha="center", style="italic")

    return save(fig, "pipeline_workflow")


# ══════════════════════════════════════════════════════════════════════
# Fig 2: Hypothesis deepening tree
# ══════════════════════════════════════════════════════════════════════

def fig_deepening_tree():
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.set_xlim(-0.5, 10)
    ax.set_ylim(-0.5, 5)
    ax.axis("off")

    # Root
    _box(ax, 2, 4, 3.5, 0.9, C_GREEN, "H-MIG01: Broadband coupling ↑",
         "d = 0.59  |  SUPPORTED", body_size=9)

    # Level 1
    _box(ax, 1.2, 2.3, 3.2, 0.9, C_BLUE, "H-NEW01: Frequency decomposition",
         "Fast ↓ / Slow ↑  |  PARTIAL", body_size=9)
    _box(ax, 5.8, 2.3, 3.2, 0.9, C_BLUE, "H-NEW17: Diurnal profile",
         "Day ↑ / Night =  |  PARTIAL", body_size=9)

    _arrow(ax, 1.5, 3.55, 1.3, 2.8, C_BLUE)
    _arrow(ax, 2.5, 3.55, 5.5, 2.8, C_BLUE)

    # Level 2
    _box(ax, 1.2, 0.5, 3.2, 0.9, C_GREEN, "H-NEW03: Excursion HR blunting",
         "d = −0.80  |  SUPPORTED", body_size=9)
    _box(ax, 5.8, 0.5, 3.2, 0.9, C_GRAY, "H-NEW06: Nocturnal → glucose",
         "ρ = −0.004  |  NULL", body_size=9)

    _arrow(ax, 1.2, 1.85, 1.2, 1.0, C_GREEN)
    _arrow(ax, 5.8, 1.85, 5.8, 1.0, C_GRAY)

    # Annotations
    ax.text(8.5, 4, "Hypothesis\nDeepening", fontsize=13, fontweight="bold",
            color=C_DARK, ha="center")
    ax.text(8.5, 3.3, "Supported findings\ngenerate more specific\nsub-hypotheses.\n\n"
            "Null findings prune\nfruitless directions.",
            fontsize=9, color="#555555", ha="center", va="top", linespacing=1.4)

    # Surprise annotation
    ax.annotate("SURPRISE:\ndirection flips\nat 2h boundary",
                xy=(1.2, 2.75), xytext=(-0.3, 3.5),
                fontsize=8, color=C_VERM, fontweight="bold", ha="center",
                arrowprops=dict(arrowstyle="-|>", color=C_VERM, lw=0.8))

    return save(fig, "deepening_tree")


# ══════════════════════════════════════════════════════════════════════
# Fig 3: 7 Quality Requirements
# ══════════════════════════════════════════════════════════════════════

def fig_quality_requirements():
    fig, ax = plt.subplots(figsize=(9, 4))
    ax.set_xlim(-0.5, 9)
    ax.set_ylim(-0.5, 7.5)
    ax.axis("off")

    reqs = [
        ("Direction", "Expected direction\n(not just 'differs')", C_BLUE),
        ("Timescale", "Specific lag\nor period range", C_SKY),
        ("Refutation", "What result\nkills this?", C_VERM),
        ("Neg. Control", "What should show\nNO effect?", C_AMBER),
        ("Confounders", "Age, site, HbA1c,\nBMI minimum", C_GREEN),
        ("Unique Enabler", "Why only\nAI-READI?", C_PINK),
        ("Effect Size", "Quantitative\nprediction", C_DARK),
    ]

    for i, (title, body, color) in enumerate(reqs):
        x = 0.5 + (i % 4) * 2.2
        y = 5.5 if i < 4 else 2.5
        _box(ax, x, y, 1.9, 1.8, color, title, body, title_size=10, body_size=8)
        # Number badge
        circle = mpatches.Circle((x - 0.75, y + 0.65), 0.2,
                                  facecolor=color, edgecolor="white", linewidth=1)
        ax.add_patch(circle)
        ax.text(x - 0.75, y + 0.65, str(i + 1), fontsize=8, fontweight="bold",
                color="white", ha="center", va="center")

    ax.text(4.25, 7.2, "Every hypothesis must specify all 7 before entering the execution queue",
            fontsize=11, fontweight="bold", color=C_DARK, ha="center")

    return save(fig, "quality_requirements")


# ══════════════════════════════════════════════════════════════════════
# Fig 4: Theory-driven vs Data-driven
# ══════════════════════════════════════════════════════════════════════

def fig_dual_mode():
    fig, axes = plt.subplots(1, 2, figsize=(11, 4.5))

    # Left: Theory-driven
    ax = axes[0]
    ax.set_xlim(-0.5, 5)
    ax.set_ylim(-0.5, 5)
    ax.axis("off")

    steps = [
        (2.5, 4.2, "Literature + Physiology", C_BLUE),
        (2.5, 3.2, "Mechanistic prediction", C_BLUE),
        (2.5, 2.2, "Designed experiment", C_BLUE),
        (2.5, 1.2, "Structured result", C_BLUE),
        (2.5, 0.2, "Verdict", C_GREEN),
    ]
    for i, (x, y, label, color) in enumerate(steps):
        rect = mpatches.FancyBboxPatch((x - 1.8, y - 0.35), 3.6, 0.7,
                                        boxstyle="round,pad=0.06", facecolor="white",
                                        edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha="center", va="center", fontsize=10, color=color, fontweight="bold")
        if i < len(steps) - 1:
            _arrow(ax, x, y - 0.4, x, steps[i+1][1] + 0.4, color, lw=1.2)

    ax.set_title("Theory-Driven", fontsize=13, fontweight="bold", color=C_BLUE, pad=10)

    # Right: Data-driven
    ax = axes[1]
    ax.set_xlim(-0.5, 5)
    ax.set_ylim(-0.5, 5)
    ax.axis("off")

    steps_d = [
        (2.5, 4.2, "Unexpected result / anomaly", C_AMBER),
        (2.5, 3.2, "Formalize into mechanism", C_AMBER),
        (2.5, 2.2, "Add direction + controls", C_AMBER),
        (2.5, 1.2, "Same execution pipeline", C_BLUE),
        (2.5, 0.2, "Same verification", C_GREEN),
    ]
    for i, (x, y, label, color) in enumerate(steps_d):
        rect = mpatches.FancyBboxPatch((x - 1.8, y - 0.35), 3.6, 0.7,
                                        boxstyle="round,pad=0.06", facecolor="white",
                                        edgecolor=color, linewidth=1.5)
        ax.add_patch(rect)
        ax.text(x, y, label, ha="center", va="center", fontsize=10, color=color, fontweight="bold")
        if i < len(steps_d) - 1:
            _arrow(ax, x, y - 0.4, x, steps_d[i+1][1] + 0.4, color, lw=1.2)

    ax.set_title("Data-Driven", fontsize=13, fontweight="bold", color=C_AMBER, pad=10)

    # Connecting annotation
    fig.text(0.5, -0.02, "Both modes produce identical structured output → same quality gate → same verification",
             ha="center", fontsize=10, color=C_DARK, style="italic")

    fig.tight_layout()
    return save(fig, "dual_mode")


# ══════════════════════════════════════════════════════════════════════
# Fig 5: Hypothesis overview heatmap-table
# ══════════════════════════════════════════════════════════════════════

def fig_hypothesis_table():
    # Load actual hypothesis data
    from hypothesis_driven.workspace import HypothesisWorkspace
    ws = HypothesisWorkspace()

    rows = []
    for h in ws.all():
        d = abs(h.results.effect_size) if h.results and h.results.effect_size else 0
        fdr = h.results.fdr_p_value if h.results else 1.0
        rows.append({
            "id": h.id,
            "category": h.category[:12],
            "status": h.status,
            "d": d,
            "sig": fdr < 0.05 if h.results else False,
            "priority": h.priority,
        })

    df = pd.DataFrame(rows).sort_values("priority", ascending=False).head(15)

    fig, ax = plt.subplots(figsize=(10, 5.5))

    status_colors = {
        "supported": C_GREEN, "completed": C_BLUE, "validated": C_SKY,
        "critiqued": C_AMBER, "inconclusive": C_GRAY, "refuted": C_VERM,
        "rejected": C_VERM, "proposed": C_SKY,
    }

    y_positions = range(len(df))
    for i, (_, row) in enumerate(df.iterrows()):
        y = len(df) - 1 - i
        color = status_colors.get(row["status"], C_GRAY)

        # ID
        ax.text(-0.1, y, row["id"], fontsize=9, fontweight="bold", va="center",
                ha="right", color=C_DARK)

        # Status badge
        badge = mpatches.FancyBboxPatch((0, y - 0.25), 1.5, 0.5,
                                         boxstyle="round,pad=0.06",
                                         facecolor=color, alpha=0.7)
        ax.add_patch(badge)
        ax.text(0.75, y, row["status"], fontsize=8, fontweight="bold",
                color="white", ha="center", va="center")

        # Category
        ax.text(2.0, y, row["category"], fontsize=9, va="center", color="#555555")

        # Effect size bar
        bar_width = min(row["d"] * 3, 3)  # scale
        bar_color = C_GREEN if row["sig"] else C_GRAY
        ax.barh(y, bar_width, left=3.5, height=0.4, color=bar_color, alpha=0.7)
        if row["d"] > 0:
            ax.text(3.5 + bar_width + 0.1, y, f"{row['d']:.2f}", fontsize=8,
                    va="center", color=bar_color, fontweight="bold")

        # Priority dot
        ax.plot(7.5, y, "o", color=color, markersize=row["priority"] * 12 + 3,
                alpha=0.6, markeredgecolor="white", markeredgewidth=0.5)

    ax.set_xlim(-1.5, 8.5)
    ax.set_ylim(-0.8, len(df) + 0.3)
    ax.axis("off")

    # Column headers
    headers = [(-0.1, "ID"), (0.75, "Status"), (2.0, "Category"), (4.5, "Effect Size |d|"), (7.5, "Priority")]
    for x, label in headers:
        ax.text(x, len(df) + 0.1, label, fontsize=10, fontweight="bold", color=C_DARK,
                ha="center" if x > 0 else "right", va="bottom")

    return save(fig, "hypothesis_table")


# ══════════════════════════════════════════════════════════════════════
# Fig 6: Comparison with SOTA systems
# ══════════════════════════════════════════════════════════════════════

def fig_comparison():
    systems = ["AI Scientist\n(2024)", "Virtual Lab\n(2024)", "SciAgents\n(2025)",
               "Agent Lab\n(2025)", "This System"]
    dims = ["Domain\nGrounding", "Hypothesis\nLifecycle", "Automated\nExecution",
            "Mechanical\nVerification", "Feedback\nLoops"]

    # Scores (0-3 scale: 0=none, 1=partial, 2=good, 3=full)
    scores = np.array([
        [0, 0, 3, 0, 1],  # AI Scientist
        [2, 1, 1, 0, 1],  # Virtual Lab
        [2, 0, 0, 0, 0],  # SciAgents
        [1, 1, 2, 0, 1],  # Agent Lab
        [3, 3, 3, 3, 3],  # Ours
    ])

    fig, ax = plt.subplots(figsize=(9, 4.5))

    colors_map = {0: "#f0f0f0", 1: "#fef3c7", 2: "#bfdbfe", 3: "#bbf7d0"}
    text_colors = {0: "#999999", 1: "#92400e", 2: "#1e40af", 3: "#166534"}
    labels_map = {0: "—", 1: "Partial", 2: "Good", 3: "Full"}

    for i in range(len(systems)):
        for j in range(len(dims)):
            s = scores[i, j]
            rect = mpatches.FancyBboxPatch((j * 1.8 + 0.05, (len(systems) - 1 - i) * 0.9 + 0.05),
                                            1.6, 0.7, boxstyle="round,pad=0.05",
                                            facecolor=colors_map[s], edgecolor="#e0e0e0", linewidth=0.5)
            ax.add_patch(rect)
            ax.text(j * 1.8 + 0.85, (len(systems) - 1 - i) * 0.9 + 0.4,
                    labels_map[s], fontsize=9, ha="center", va="center",
                    color=text_colors[s], fontweight="bold" if s == 3 else "normal")

    # Row labels
    for i, sys in enumerate(systems):
        weight = "bold" if i == 4 else "normal"
        color = C_GREEN if i == 4 else C_DARK
        ax.text(-0.15, (len(systems) - 1 - i) * 0.9 + 0.4, sys,
                fontsize=10, ha="right", va="center", fontweight=weight, color=color)

    # Column labels
    for j, dim in enumerate(dims):
        ax.text(j * 1.8 + 0.85, len(systems) * 0.9 + 0.15, dim,
                fontsize=9, ha="center", va="bottom", fontweight="bold", color=C_DARK)

    ax.set_xlim(-2.5, len(dims) * 1.8 + 0.2)
    ax.set_ylim(-0.3, len(systems) * 0.9 + 0.6)
    ax.axis("off")

    return save(fig, "sota_comparison")


# ══════════════════════════════════════════════════════════════════════
# Fig 7: Verifier 9 checks
# ══════════════════════════════════════════════════════════════════════

def fig_verifier_checks():
    checks = [
        ("Covariate\nadjustment", "age + HbA1c\n+ BMI + site"),
        ("Negative\ncontrol", "Random windows\n= no effect"),
        ("Dose-\nresponse", "Monotonic\nacross groups"),
        ("Bootstrap\nCI", "95% CI\nexcludes zero"),
        ("Sample\nsize", "N ≥ 100\nper group"),
        ("Effect\nsize", "|d| ≥ 0.2"),
        ("FDR\nsignificance", "FDR p < 0.05"),
        ("Site\nbias", "Site in\ncovariates"),
        ("Sensitivity\nsurvival", "≥ 67%\nchecks pass"),
    ]

    fig, ax = plt.subplots(figsize=(10, 3.5))
    ax.set_xlim(-0.5, 10)
    ax.set_ylim(-0.5, 3.5)
    ax.axis("off")

    for i, (title, detail) in enumerate(checks):
        x = (i % 5) * 2.0 + 0.5
        y = 2.2 if i < 5 else 0.3

        # Circle with number
        circle = mpatches.Circle((x, y + 0.45), 0.25, facecolor=C_GREEN, edgecolor="white", linewidth=1.5)
        ax.add_patch(circle)
        ax.text(x, y + 0.45, str(i + 1), fontsize=9, fontweight="bold",
                color="white", ha="center", va="center")

        # Label
        ax.text(x, y, title, fontsize=9, fontweight="bold", ha="center", va="top", color=C_DARK)
        ax.text(x, y - 0.55, detail, fontsize=7, ha="center", va="top", color="#666666", linespacing=1.2)

    # Decision rule
    ax.text(5, -0.3, "≥ 7/9 pass + PASS → supported   |   ≤ 3/9 or FAIL → refuted   |   else → inconclusive",
            fontsize=10, ha="center", color=C_DARK, fontweight="bold",
            bbox=dict(boxstyle="round,pad=0.3", facecolor="#f0f0f0", edgecolor="#dddddd"))

    ax.text(5, 3.3, "9 Mechanical Robustness Checks (rule-based, not LLM)",
            fontsize=12, fontweight="bold", color=C_DARK, ha="center")

    return save(fig, "verifier_checks")


# ══════════════════════════════════════════════════════════════════════
# PPTX Assembly
# ══════════════════════════════════════════════════════════════════════

BG = RGBColor(0x1a, 0x1a, 0x2e)
SB = RGBColor(0x55, 0x55, 0x55)
AC = RGBColor(0x00, 0x72, 0xB2)
WH = RGBColor(0xff, 0xff, 0xff)


def _title(slide, title, subtitle=""):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.85))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = BG
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle
        p2.font.size = Pt(11); p2.font.color.rgb = SB


def img_slide(prs, title, img, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, title, subtitle)
    if img and os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.5), Inches(1.05), Inches(12))


def main():
    print("Generating system design figures...")
    f1 = fig_pipeline()
    f2 = fig_deepening_tree()
    f3 = fig_quality_requirements()
    f4 = fig_dual_mode()
    f5 = fig_hypothesis_table()
    f6 = fig_comparison()
    f7 = fig_verifier_checks()
    print(f"  7 figures → {FIGDIR}/")

    print("Assembling PPTX...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tx = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = "Hypothesis-Driven Scientific Discovery System"
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = BG
    p2 = tf.add_paragraph(); p2.text = "\nSystematic Design for Agentic Multi-Modal Physiological Analysis"
    p2.font.size = Pt(16); p2.font.color.rgb = SB; p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph(); p3.text = "\nMechanism first, measurement second"
    p3.font.size = Pt(13); p3.font.color.rgb = AC; p3.alignment = PP_ALIGN.CENTER; p3.font.italic = True

    img_slide(prs, "Pipeline Workflow", f1,
              "Proposer → Critic → Workspace → Executor → Verifier with feedback loops")
    img_slide(prs, "Dual-Mode Hypothesis Generation", f4,
              "Theory-driven and data-driven hypotheses enter the same quality gate and verification pipeline")
    img_slide(prs, "7 Quality Requirements", f3,
              "Every hypothesis must specify direction, timescale, refutation criterion, negative control, confounders, unique enabler, and effect size")
    img_slide(prs, "Hypothesis Deepening: From Broad to Specific", f2,
              "Supported findings generate sub-hypotheses. Null findings prune fruitless directions. Surprises seed data-driven hypotheses.")
    img_slide(prs, "Hypothesis Overview", f5,
              "26 hypotheses across 10 categories. Status, effect sizes, and priority shown.")
    img_slide(prs, "9 Mechanical Robustness Checks", f7,
              "Rule-based verification takes precedence over LLM judgment. Prevents hallucination from overriding evidence.")
    img_slide(prs, "Comparison with State-of-the-Art Systems", f6,
              "Our system uniquely combines domain grounding, structured lifecycle, automated execution, mechanical verification, and feedback loops.")

    out = Path("docs/hypothesis_system_deck.pptx")
    prs.save(out)
    print(f"Saved: {out} ({out.stat().st_size / 1024:.0f} KB)")
    print(f"Editable vectors: {FIGDIR}/*.svg and *.pdf")


if __name__ == "__main__":
    main()
