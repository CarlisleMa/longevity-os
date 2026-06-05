#!/usr/bin/env python3
"""Build a fresh 3-slide Longevity OS intro deck.

The deck is intentionally independent from the existing AI-READI slide decks.
It creates new scientific figures and a short companion speaker-notes file.
"""

from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "decks" / "longevity_os_intro"
FIG_DIR = OUT_DIR / "figures"
OUT_PPTX = OUT_DIR / "longevity_os_intro_3slides.pptx"
OUT_NOTES = OUT_DIR / "longevity_os_intro_speaker_notes.md"


# Colorblind-safe, conference-friendly palette.
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(24, 29, 36)
MUTED = RGBColor(88, 96, 108)
LIGHT = RGBColor(232, 235, 240)
FAINT = RGBColor(247, 248, 250)
BLUE = RGBColor(0, 114, 178)
SKY = RGBColor(86, 180, 233)
GREEN = RGBColor(0, 158, 115)
ORANGE = RGBColor(230, 159, 0)
VERMILION = RGBColor(213, 94, 0)
PURPLE = RGBColor(117, 112, 179)

C_INK = "#181d24"
C_MUTED = "#58606c"
C_LIGHT = "#e8ebf0"
C_BLUE = "#0072B2"
C_SKY = "#56B4E9"
C_GREEN = "#009E73"
C_ORANGE = "#E69F00"
C_VERMILION = "#D55E00"
C_PURPLE = "#7570B3"


plt.rcParams.update(
    {
        "figure.facecolor": "white",
        "axes.facecolor": "white",
        "savefig.facecolor": "white",
        "font.family": "DejaVu Sans",
        "font.size": 10,
        "axes.titlesize": 12,
        "axes.labelsize": 10,
        "xtick.labelsize": 9,
        "ytick.labelsize": 9,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "axes.edgecolor": "#9aa3af",
        "axes.linewidth": 0.8,
        "xtick.color": C_MUTED,
        "ytick.color": C_MUTED,
        "axes.labelcolor": C_MUTED,
        "text.color": C_INK,
    }
)


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    FIG_DIR.mkdir(parents=True, exist_ok=True)


def save_figure(fig: plt.Figure, stem: str) -> Path:
    png = FIG_DIR / f"{stem}.png"
    svg = FIG_DIR / f"{stem}.svg"
    pdf = FIG_DIR / f"{stem}.pdf"
    fig.savefig(png, dpi=300, bbox_inches="tight", pad_inches=0.08)
    fig.savefig(svg, bbox_inches="tight", pad_inches=0.08)
    fig.savefig(pdf, bbox_inches="tight", pad_inches=0.08)
    plt.close(fig)
    return png


def rgb_to_hex(color: RGBColor) -> str:
    return f"#{color[0]:02x}{color[1]:02x}{color[2]:02x}"


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if line_spacing is not None:
        p.line_spacing = line_spacing
    return box


def add_multiline(
    slide,
    lines: list[tuple[str, int, RGBColor, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    spacing_after: int = 3,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(spacing_after)
    return box


def add_rule(slide, x: float, y: float, w: float, color: RGBColor = BLUE) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.035)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_header(slide, title: str, subtitle: str, number: str) -> None:
    set_background(slide)
    add_text(slide, "Longevity OS", 0.52, 0.22, 1.8, 0.28, 9, MUTED, True)
    add_text(slide, number, 12.25, 0.2, 0.45, 0.3, 9, MUTED, False, PP_ALIGN.RIGHT)
    add_text(slide, title, 0.52, 0.5, 8.8, 0.45, 23, INK, True)
    add_text(slide, subtitle, 0.52, 0.96, 9.7, 0.32, 10, MUTED)
    add_rule(slide, 0.52, 1.25, 1.35, BLUE)


def add_panel(slide, x: float, y: float, w: float, h: float, title: str, color: RGBColor):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = FAINT
    rect.line.color.rgb = LIGHT
    rect.line.width = Pt(0.7)
    add_text(slide, title, x + 0.18, y + 0.14, w - 0.36, 0.28, 9, color, True)
    return rect


def add_caption(slide, text: str) -> None:
    add_text(slide, text, 0.55, 7.12, 12.25, 0.22, 6.8, MUTED)


def fig_reserve_model() -> Path:
    fig = plt.figure(figsize=(6.9, 3.9))
    gs = fig.add_gridspec(1, 2, width_ratios=[1.12, 1.0], wspace=0.35)

    ax = fig.add_subplot(gs[0, 0])
    age = np.linspace(35, 90, 300)
    reserve_a = 1 / (1 + np.exp((age - 76) / 7.5))
    reserve_b = 1 / (1 + np.exp((age - 61) / 7.5))
    ax.plot(age, reserve_a, color=C_BLUE, lw=2.5, label="Resilient trajectory")
    ax.plot(age, reserve_b, color=C_VERMILION, lw=2.5, label="Vulnerable trajectory")
    ax.axvline(60, color="#333333", lw=1, ls="--")
    ax.scatter([60, 60], [np.interp(60, age, reserve_a), np.interp(60, age, reserve_b)],
               s=62, color=[C_BLUE, C_VERMILION], edgecolor="white", linewidth=1.0, zorder=5)
    ax.text(61.0, 0.84, "same age", fontsize=7.5, color=C_INK)
    ax.text(72, 0.23, "clinical threshold", fontsize=8, color=C_MUTED)
    ax.axhline(0.32, color="#9aa3af", lw=1, ls=":")
    ax.set_xlim(38, 88)
    ax.set_ylim(0.02, 1.02)
    ax.set_xlabel("Chronological age")
    ax.set_ylabel("Physiological reserve")
    ax.text(0.0, 1.04, "A. same age, different reserve", transform=ax.transAxes,
            fontsize=8.5, fontweight="bold", ha="left")
    ax.legend(frameon=False, fontsize=8, loc="lower left")

    ax2 = fig.add_subplot(gs[0, 1])
    ax2.axis("off")
    center = np.array([0.5, 0.52])
    diseases = [
        ("T2D", 0.50, 0.89, C_ORANGE),
        ("CVD", 0.87, 0.68, C_VERMILION),
        ("CKD", 0.77, 0.23, C_SKY),
        ("frailty", 0.23, 0.23, C_PURPLE),
        ("cognitive\nrisk", 0.12, 0.68, C_GREEN),
    ]
    for label, x, y, c in diseases:
        ax2.plot([center[0], x], [center[1], y], color="#c8ced8", lw=1.6, zorder=1)
        circ = patches.Circle((x, y), 0.1, facecolor="white", edgecolor=c, lw=2.0, zorder=3)
        ax2.add_patch(circ)
        ax2.text(x, y, label, ha="center", va="center", fontsize=8, color=C_INK, zorder=4)
    hub = patches.Circle(center, 0.16, facecolor=C_INK, edgecolor=C_INK, lw=1.0, zorder=5)
    ax2.add_patch(hub)
    ax2.text(center[0], center[1], "aging\nstate", ha="center", va="center",
             fontsize=9, color="white", zorder=6)
    ax2.text(0.0, 0.04, "A latent state that changes disease\nsusceptibility across organs",
             fontsize=8.5, color=C_MUTED, ha="left")
    ax2.text(0.0, 1.04, "B. shared substrate for disease", transform=ax2.transAxes,
             fontsize=8.5, fontweight="bold", ha="left")
    ax2.set_xlim(0, 1)
    ax2.set_ylim(0, 1)
    fig.subplots_adjust(top=0.88, bottom=0.16, left=0.08, right=0.98)
    return save_figure(fig, "slide1_aging_system_state")


def fig_measurement_landscape() -> Path:
    fig, ax = plt.subplots(figsize=(8.0, 4.25))
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)

    rows = [
        ("Clinical biomarkers", "labs, vitals, function", "risk factors / composite age", C_BLUE),
        ("Phenotypic clocks", "KDM, PhenoAge, frailty", "mortality and healthspan risk", C_GREEN),
        ("Molecular clocks", "DNAm, proteins, metabolites", "age, lifespan, pace of aging", C_ORANGE),
        ("Organ clocks", "retina, ECG, brain, imaging", "organ-specific age gap", C_PURPLE),
        ("Digital physiology", "CGM, HR, sleep, activity", "dynamics and regulation", C_VERMILION),
    ]

    x0 = 0.04
    widths = [0.25, 0.28, 0.31]
    headers = ["Measurement family", "Input signal", "Typical output"]
    xs = [x0, x0 + widths[0] + 0.025, x0 + widths[0] + widths[1] + 0.05]
    for x, w, h in zip(xs, widths, headers):
        ax.text(x, 0.94, h, fontsize=9, color=C_MUTED, fontweight="bold", va="center")
        ax.plot([x, x + w], [0.905, 0.905], color="#c8ced8", lw=1.2)

    y_top = 0.82
    row_h = 0.145
    for i, (family, signal, output, color) in enumerate(rows):
        y = y_top - i * row_h
        bg = "#fbfcfd" if i % 2 == 0 else "#f4f6f8"
        rect = patches.FancyBboxPatch((0.025, y - 0.055), 0.93, 0.095,
                                      boxstyle="round,pad=0.012,rounding_size=0.012",
                                      facecolor=bg, edgecolor="#e4e8ee", lw=0.8)
        ax.add_patch(rect)
        ax.add_patch(patches.Circle((0.055, y - 0.008), 0.018, facecolor=color, edgecolor="none"))
        ax.text(xs[0] + 0.035, y, family, fontsize=9.5, color=C_INK, fontweight="bold", va="center")
        ax.text(xs[1], y, signal, fontsize=9, color=C_INK, va="center")
        ax.text(xs[2], y, output, fontsize=9, color=C_INK, va="center")

    ax.annotate("", xy=(0.955, 0.18), xytext=(0.955, 0.84),
                arrowprops=dict(arrowstyle="-|>", lw=1.3, color=C_MUTED))
    ax.text(0.976, 0.51, "more multimodal,\nless standardized", fontsize=8,
            rotation=90, color=C_MUTED, ha="center", va="center")

    ax.text(0.04, 0.075,
            "Field trajectory: from single risk markers -> supervised age/risk clocks -> longitudinal pace and organ-specific models.",
            fontsize=8.5, color=C_MUTED)
    return save_figure(fig, "slide2_measurement_landscape")


def fig_missing_layer() -> Path:
    fig = plt.figure(figsize=(8.2, 3.9))
    gs = fig.add_gridspec(1, 2, width_ratios=[0.9, 1.25], wspace=0.22)

    ax = fig.add_subplot(gs[0, 0])
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)

    # Snapshot clock cartoon.
    ax.add_patch(patches.FancyBboxPatch((0.08, 0.67), 0.28, 0.17, boxstyle="round,pad=0.02",
                                        facecolor="#f4f6f8", edgecolor="#cbd2dc", lw=1.0))
    ax.text(0.22, 0.755, "x(t0)", ha="center", va="center", fontsize=11, fontweight="bold")
    ax.text(0.22, 0.64, "snapshot", ha="center", va="top", fontsize=7.5, color=C_MUTED)
    ax.annotate("", xy=(0.56, 0.755), xytext=(0.38, 0.755),
                arrowprops=dict(arrowstyle="-|>", lw=1.4, color=C_MUTED))
    ax.add_patch(patches.FancyBboxPatch((0.58, 0.67), 0.32, 0.17, boxstyle="round,pad=0.02",
                                        facecolor="white", edgecolor=C_BLUE, lw=1.8))
    ax.text(0.74, 0.765, "age gap", ha="center", va="center", fontsize=10.5, color=C_BLUE, fontweight="bold")
    ax.text(0.74, 0.64, "one compressed scalar", ha="center", va="top", fontsize=7.5, color=C_MUTED)
    ax.text(0.02, 1.04, "A. snapshot-clock abstraction", transform=ax.transAxes,
            fontsize=8.5, fontweight="bold", ha="left")

    ax.add_patch(patches.FancyBboxPatch((0.08, 0.34), 0.82, 0.17, boxstyle="round,pad=0.02",
                                        facecolor="#fbfcfd", edgecolor="#e4e8ee", lw=0.8))
    ax.text(0.12, 0.45, "Captures", fontsize=8.2, fontweight="bold", color=C_INK)
    ax.text(0.34, 0.45, "levels, burden, risk proxy", fontsize=8.0, color=C_MUTED)
    ax.add_patch(patches.FancyBboxPatch((0.08, 0.12), 0.82, 0.17, boxstyle="round,pad=0.02",
                                        facecolor="#fffaf5", edgecolor="#f1d7bd", lw=0.8))
    ax.text(0.12, 0.23, "Loses", fontsize=8.2, fontweight="bold", color=C_VERMILION)
    ax.text(0.34, 0.23, "response, lag, phase, feedback", fontsize=8.0, color=C_MUTED)

    ax2 = fig.add_subplot(gs[0, 1])
    ax2.axis("off")
    ax2.set_xlim(0, 1)
    ax2.set_ylim(0, 1)
    nodes = {
        "glucose": (0.50, 0.86, C_ORANGE),
        "autonomic\nHR": (0.82, 0.62, C_BLUE),
        "activity": (0.70, 0.22, C_GREEN),
        "sleep": (0.30, 0.22, C_PURPLE),
        "environment": (0.18, 0.62, C_SKY),
        "clinical\nretina ECG": (0.50, 0.51, C_INK),
    }
    edges = [
        ("glucose", "autonomic\nHR", "lag"),
        ("autonomic\nHR", "activity", "response"),
        ("activity", "sleep", "recovery"),
        ("sleep", "glucose", "next-day"),
        ("environment", "glucose", "exposure"),
        ("environment", "autonomic\nHR", "stress"),
        ("clinical\nretina ECG", "glucose", "structure"),
        ("clinical\nretina ECG", "autonomic\nHR", "organ state"),
        ("clinical\nretina ECG", "sleep", "reserve"),
        ("clinical\nretina ECG", "activity", "function"),
    ]
    for a, b, label in edges:
        x1, y1, _ = nodes[a]
        x2, y2, _ = nodes[b]
        ax2.plot([x1, x2], [y1, y2], color="#c8ced8", lw=1.4, zorder=1)
        mx, my = (x1 + x2) / 2, (y1 + y2) / 2
        if label in {"lag", "response", "recovery", "next-day", "exposure", "stress"}:
            ax2.text(mx, my, label, fontsize=6.8, color=C_MUTED, ha="center", va="center",
                     bbox=dict(boxstyle="round,pad=0.12", fc="white", ec="none", alpha=0.9))
    for name, (x, y, color) in nodes.items():
        r = 0.085 if name != "clinical\nretina ECG" else 0.105
        ax2.add_patch(patches.Circle((x, y), r, facecolor="white", edgecolor=color, lw=2.0, zorder=3))
        ax2.text(x, y, name, ha="center", va="center", fontsize=8.2, color=C_INK, zorder=4)
    ax2.text(0.04, 0.02,
             "Needed representation: system state + coupling graph + response to daily-life perturbations.",
             fontsize=8.5, color=C_MUTED)
    ax2.text(0.0, 1.04, "B. coordination graph", transform=ax2.transAxes,
             fontsize=8.5, fontweight="bold", ha="left")
    fig.subplots_adjust(top=0.88, bottom=0.12, left=0.04, right=0.98)
    return save_figure(fig, "slide3_missing_coordination")


def build_deck(fig1: Path, fig2: Path, fig3: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1.
    slide = prs.slides.add_slide(blank)
    add_header(
        slide,
        "Aging is a latent systems variable",
        "The same chronological age can encode different vulnerability, reserve, and disease trajectory.",
        "01",
    )
    add_multiline(
        slide,
        [
            ("Problem framing", 10, BLUE, True),
            ("Chronological age is the label we observe, not the biology we want to measure.", 19, INK, True),
            ("Aging manifests as declining reserve, impaired repair, chronic inflammation, weaker stress response, and rising multi-disease susceptibility.", 11, MUTED, False),
            ("The scientific target is therefore not just age prediction. It is estimating the hidden physiological state that explains why people of the same age diverge.", 11, MUTED, False),
        ],
        0.65,
        1.62,
        4.0,
        3.0,
    )
    add_panel(slide, 0.65, 4.68, 4.0, 1.45, "Opening claim", BLUE)
    add_text(
        slide,
        "Longevity begins when aging becomes measurable as a modifiable system state.",
        0.88,
        5.02,
        3.55,
        0.62,
        17,
        INK,
        True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    slide.shapes.add_picture(str(fig1), Inches(5.05), Inches(1.55), width=Inches(7.65))
    add_caption(
        slide,
        "Conceptual anchor: the 2023 hallmarks framework defines aging as multi-layered loss of integrity, not a single observable measurement.",
    )

    # Slide 2.
    slide = prs.slides.add_slide(blank)
    add_header(
        slide,
        "Biological age measurement has moved from markers to clocks",
        "The field now estimates age, mortality risk, organ age, and pace of aging from increasingly high-dimensional data.",
        "02",
    )
    slide.shapes.add_picture(str(fig2), Inches(0.62), Inches(1.55), width=Inches(7.7))
    add_panel(slide, 8.65, 1.55, 3.9, 0.92, "What changed", GREEN)
    add_text(
        slide,
        "Aging measurement shifted from individual biomarkers to supervised representations trained on age, mortality, disease, or longitudinal decline.",
        8.88,
        1.92,
        3.42,
        0.32,
        9.5,
        MUTED,
    )
    add_panel(slide, 8.65, 2.75, 3.9, 1.38, "Canonical examples", ORANGE)
    add_text(
        slide,
        "Horvath DNAm clock: chronological age across tissues\nPhenoAge: clinical phenotypic risk mapped to methylation\nGrimAge: mortality-linked DNAm surrogates\nDunedinPACE: longitudinal pace of aging",
        8.88,
        3.08,
        3.42,
        0.78,
        8.7,
        MUTED,
    )
    add_panel(slide, 8.65, 4.42, 3.9, 1.28, "Core abstraction", BLUE)
    add_text(
        slide,
        "clock gap = predicted biological age - chronological age",
        8.88,
        4.79,
        3.42,
        0.24,
        10,
        INK,
        True,
    )
    add_text(
        slide,
        "Useful, compact, and predictive; but often optimized to reconstruct labels rather than physiology.",
        8.88,
        5.12,
        3.42,
        0.34,
        9,
        MUTED,
    )
    add_caption(
        slide,
        "Representative sources: Horvath 2013; Levine et al. 2018; Lu et al. 2019; Belsky et al. 2022.",
    )

    # Slide 3.
    slide = prs.slides.add_slide(blank)
    add_header(
        slide,
        "The deficiency: clocks are snapshots of an interacting system",
        "Longevity depends on regulation, coupling, and resilience under daily-life perturbations.",
        "03",
    )
    slide.shapes.add_picture(str(fig3), Inches(0.62), Inches(1.55), width=Inches(8.05))
    add_panel(slide, 9.0, 1.55, 3.65, 1.1, "Missing axes", VERMILION)
    add_text(
        slide,
        "1. Dynamics: how fast systems respond\n2. Coupling: which systems move together\n3. Context: sleep, activity, environment\n4. Resilience: recovery after perturbation",
        9.24,
        1.89,
        3.16,
        0.64,
        8.8,
        MUTED,
    )
    add_panel(slide, 9.0, 2.95, 3.65, 1.22, "Why diabetes is a strong aging lens", BLUE)
    add_text(
        slide,
        "Type 2 diabetes is age-associated, multisystem, and dynamic. It perturbs metabolism, autonomic tone, sleep, activity, vascular structure, and inflammatory burden.",
        9.24,
        3.31,
        3.16,
        0.58,
        9.1,
        MUTED,
    )
    add_panel(slide, 9.0, 4.48, 3.65, 1.22, "Longevity OS thesis", GREEN)
    add_text(
        slide,
        "Measure aging as a multimodal state plus a cross-system coordination graph, not as one scalar clock.",
        9.24,
        4.88,
        3.16,
        0.46,
        12,
        INK,
        True,
    )
    add_caption(
        slide,
        "This motivates AI-READI: clinical labs + retina + ECG + CGM + wearable sleep/activity/HR + environment in the same participants.",
    )

    prs.save(OUT_PPTX)


def write_notes() -> None:
    notes = """# Longevity OS intro slides - speaker notes

## Slide 1 - Aging is a latent systems variable

The opening move is to separate chronological age from biological state. Chronological age is the observable label, but the scientific object is hidden: reserve, repair capacity, inflammatory tone, metabolic flexibility, vascular integrity, tissue structure, and stress response. The key motivation is that two people can be the same age while carrying very different vulnerability. That gives Longevity OS a purpose: infer the system state that explains divergence.

## Slide 2 - From markers to clocks

This slide gives credit to the aging-clock field. The point is not to dismiss clocks. Clinical biomarkers, phenotypic clocks, molecular clocks, organ clocks, and digital physiology all made aging increasingly measurable. The transition to make in speech is that most clocks are supervised mappings from high-dimensional data to labels such as chronological age, mortality, or longitudinal decline. That is powerful, but it is also a compression.

## Slide 3 - The deficiency

The deficiency is that aging is interactive and dynamic, while most clocks are static. They capture levels and burden, but often miss lag, phase, feedback, stress response, recovery, and context. This sets up the project: AI-READI is valuable because it measures clinical state, retina, ECG, CGM, wearable physiology, sleep/activity, and environment in the same people. The thesis is that longevity should be modeled as a multimodal state plus a coordination graph.

## Compact transition into slide 4

Therefore, Longevity OS is not another single clock. It is an operating system for multimodal aging measurement: static organ age, dynamic physiology, cross-system coupling, and hypothesis-driven scientific reasoning.
"""
    OUT_NOTES.write_text(notes)


def main() -> None:
    ensure_dirs()
    fig1 = fig_reserve_model()
    fig2 = fig_measurement_landscape()
    fig3 = fig_missing_layer()
    build_deck(fig1, fig2, fig3)
    write_notes()
    print(f"Wrote {OUT_PPTX}")
    print(f"Wrote {OUT_NOTES}")
    print(f"Wrote figures to {FIG_DIR}")


if __name__ == "__main__":
    main()
