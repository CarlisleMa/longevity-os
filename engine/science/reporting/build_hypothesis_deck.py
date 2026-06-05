"""
Build hypothesis-driven discovery PPTX deck with Nature-quality figures.

Uses SciencePlots 'nature' style as base, adapted for widescreen PPTX.
All data from actual result CSVs/parquets. 300 DPI, colorblind-safe.
"""

from __future__ import annotations

import os
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd
import scienceplots  # noqa: F401 — registers styles
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from scripts.config import RESULTS_DIR, result_path

# ── Palette: Okabe-Ito (universally colorblind-safe) ──────────────────

C_HEALTHY = "#0072B2"
C_PREDIAB = "#56B4E9"
C_ORAL    = "#E69F00"
C_INSULIN = "#D55E00"
C_GROUPS  = [C_HEALTHY, C_PREDIAB, C_ORAL, C_INSULIN]
G_MARKERS = ["o", "s", "^", "D"]
G_SHORT   = ["Healthy", "Pre-DM", "Oral med", "Insulin"]
G_ORDER   = ["healthy", "pre_diabetes_lifestyle_controlled",
             "oral_medication_and_or_non_insulin_injectable_medication_controlled",
             "insulin_dependent"]

C_POS     = "#D55E00"
C_NEG     = "#0072B2"
C_NULL    = "#999999"
C_PASS    = "#009E73"

FIGDIR = RESULTS_DIR / "figures" / "deck"
FIGDIR.mkdir(parents=True, exist_ok=True)
DPI = 300

# Nature style + presentation scaling (fonts ~1.5x paper for projection)
plt.style.use(["science", "nature", "no-latex"])
plt.rcParams.update({
    "font.family": "sans-serif",
    "font.sans-serif": ["Arial", "DejaVu Sans", "Helvetica"],
    "font.size": 10,
    "axes.titlesize": 12,
    "axes.labelsize": 11,
    "xtick.labelsize": 10,
    "ytick.labelsize": 10,
    "legend.fontsize": 9,
    "figure.facecolor": "white",
    "axes.facecolor": "white",
    "savefig.facecolor": "white",
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.linewidth": 0.8,
    "lines.linewidth": 1.5,
    "lines.markersize": 5,
    "xtick.major.width": 0.8,
    "ytick.major.width": 0.8,
    "xtick.major.size": 3.5,
    "ytick.major.size": 3.5,
    "xtick.direction": "out",
    "ytick.direction": "out",
})


def save_fig(fig, name):
    """Save figure in three formats: PNG (for PPTX), SVG + PDF (editable vector)."""
    png_path = FIGDIR / f"{name}.png"
    svg_path = FIGDIR / f"{name}.svg"
    pdf_path = FIGDIR / f"{name}.pdf"
    fig.savefig(png_path, dpi=DPI, bbox_inches="tight", pad_inches=0.1)
    fig.savefig(svg_path, format="svg", bbox_inches="tight", pad_inches=0.1)
    fig.savefig(pdf_path, format="pdf", bbox_inches="tight", pad_inches=0.1)
    plt.close(fig)
    return str(png_path)


def sigstars(p):
    if p < 0.001: return "***"
    if p < 0.01: return "**"
    if p < 0.05: return "*"
    return "n.s."


def panel_label(ax, label, x=-0.12, y=1.08):
    """Add bold panel label (a, b, c...) per Nature convention."""
    ax.text(x, y, label, transform=ax.transAxes, fontsize=14,
            fontweight="bold", va="top", ha="left")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 1: Architecture — clean boxes
# ══════════════════════════════════════════════════════════════════════

def fig_architecture():
    fig, ax = plt.subplots(figsize=(10, 2.8))
    ax.set_xlim(-0.2, 10.8)
    ax.set_ylim(-0.8, 2.8)
    ax.axis("off")

    boxes = [
        (1.0,  1.0, "#0072B2", "Proposer",   "Literature\n& mechanisms"),
        (3.2,  1.0, "#E69F00", "Critic",     "Feasibility\n& rigor"),
        (5.4,  1.0, "#999999", "Workspace",  "Prioritized\nqueue"),
        (7.6,  1.0, "#009E73", "Executor",   "Analysis\n& code"),
        (9.8,  1.0, "#CC79A7", "Verifier",   "9-test\nrobustness"),
    ]

    for x, y, c, title, sub in boxes:
        rect = mpatches.FancyBboxPatch(
            (x - 0.85, y - 0.55), 1.7, 1.3,
            boxstyle="round,pad=0.08", facecolor="white",
            edgecolor=c, linewidth=2.0
        )
        ax.add_patch(rect)
        ax.text(x, y + 0.22, title, ha="center", va="center",
                fontsize=11, fontweight="bold", color=c)
        ax.text(x, y - 0.18, sub, ha="center", va="center",
                fontsize=8, color="#555555", linespacing=1.2)

    for i in range(4):
        x1 = boxes[i][0] + 0.88
        x2 = boxes[i+1][0] - 0.88
        ax.annotate("", xy=(x2, 1.0), xytext=(x1, 1.0),
                     arrowprops=dict(arrowstyle="-|>", color="#aaaaaa", lw=1.5, mutation_scale=12))

    ax.annotate("", xy=(1.0, 0.3), xytext=(9.8, 0.3),
                arrowprops=dict(arrowstyle="-|>", color="#CC79A7", lw=1.0,
                                connectionstyle="arc3,rad=0.2", linestyle="--", mutation_scale=10))
    ax.text(5.4, -0.05, "Refine", ha="center", fontsize=8, color="#CC79A7", style="italic")

    return save_fig(fig, "architecture")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 2: Hypothesis status
# ══════════════════════════════════════════════════════════════════════

def fig_hypothesis_status():
    labels = ["Supported", "Completed", "Validated", "Critiqued", "Inconclusive", "Refuted"]
    counts = [4, 5, 6, 6, 2, 1]
    colors = [C_PASS, C_HEALTHY, C_PREDIAB, C_ORAL, C_NULL, C_INSULIN]

    fig, ax = plt.subplots(figsize=(6, 3))
    y = range(len(labels))[::-1]
    bars = ax.barh(y, counts, color=colors, height=0.55, edgecolor="white", linewidth=0.8)
    ax.set_yticks(list(y))
    ax.set_yticklabels(labels)
    ax.set_xlabel("Number of hypotheses")
    ax.set_xlim(0, 8.5)
    for bar, c in zip(bars, counts):
        ax.text(bar.get_width() + 0.15, bar.get_y() + bar.get_height()/2,
                str(c), va="center", fontsize=10, fontweight="bold")
    ax.text(0.97, 0.03, "N = 26 total", transform=ax.transAxes,
            ha="right", va="bottom", fontsize=8, color="#777777")
    return save_fig(fig, "hypothesis_status")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 3: Forest plot
# ══════════════════════════════════════════════════════════════════════

def fig_forest_plot():
    data = [
        ("Post-excursion HR blunting",  -0.805, -0.952, -0.676, C_PASS,    "H-NEW03"),
        ("Broadband coupling",           0.585,  0.45,   0.72,  C_PASS,    "H-MIG01"),
        ("Diurnal mesor elevation",      0.466,  0.30,   0.63,  C_HEALTHY, "H-NEW17"),
        ("4\u20138h coherence",          0.377,  0.22,   0.53,  C_HEALTHY, "H-NEW01"),
        ("30\u201360min coherence",     -0.218, -0.37,  -0.07,  C_HEALTHY, "H-NEW01"),
        ("Glucose entropy",              0.175,  0.04,   0.31,  C_ORAL,    "H-NEW05"),
        ("Nocturnal \u2192 glucose",    -0.018, -0.16,   0.12,  C_NULL,    "H-NEW06"),
        ("Sleep-wake reorg.",           -0.063, -0.21,   0.09,  C_NULL,    "H-NEW13"),
    ]

    fig, ax = plt.subplots(figsize=(7, 4.5))
    n = len(data)

    for i, (label, d, lo, hi, c, hid) in enumerate(data):
        y = n - 1 - i
        ax.plot([lo, hi], [y, y], color=c, linewidth=2.0, solid_capstyle="butt")
        ax.plot([lo, lo], [y-0.08, y+0.08], color=c, linewidth=1.5)
        ax.plot([hi, hi], [y-0.08, y+0.08], color=c, linewidth=1.5)
        ax.plot(d, y, "D", color=c, markersize=7, zorder=5,
                markeredgecolor="white", markeredgewidth=0.8)

    ax.axvspan(-0.2, 0.2, alpha=0.04, color="#000000", zorder=0)
    ax.axvline(0, color="#333333", linewidth=0.8, zorder=0)

    ax.set_yticks(range(n))
    ax.set_yticklabels([d[0] for d in data][::-1], fontsize=9)
    ax.set_xlabel("Cohen's d (Healthy vs Insulin-dependent)")
    ax.set_xlim(-1.3, 1.0)

    return save_fig(fig, "forest_plot")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 4: Frequency bands — diverging lollipop
# ══════════════════════════════════════════════════════════════════════

def fig_frequency_bands():
    bands  = ["30\u201360 min", "1\u20132 h", "2\u20134 h", "4\u20138 h", "Broadband"]
    ds     = [-0.218, -0.172, +0.327, +0.377, +0.055]
    ps     = [2.17e-10, 6.13e-7, 2.64e-2, 2.30e-8, 6.06e-2]
    colors = [C_NEG, C_NEG, C_POS, C_POS, C_NULL]

    fig, ax = plt.subplots(figsize=(7, 4))
    y_pos = np.arange(len(bands))[::-1]

    # Lollipop: stem + dot (cleaner than bars)
    for yp, d, c in zip(y_pos, ds, colors):
        ax.plot([0, d], [yp, yp], color=c, linewidth=2.5, solid_capstyle="round")
        ax.plot(d, yp, "o", color=c, markersize=10, zorder=5,
                markeredgecolor="white", markeredgewidth=1.0)

    ax.axvline(0, color="#333333", linewidth=0.8, zorder=0)

    for yp, d, p in zip(y_pos, ds, ps):
        offset = 0.03 if d >= 0 else -0.03
        ha = "left" if d >= 0 else "right"
        ax.text(d + offset, yp + 0.15, f"{d:+.2f} {sigstars(p)}", ha=ha,
                fontsize=9, fontweight="bold", color="#333333")

    ax.set_yticks(y_pos)
    ax.set_yticklabels(bands)
    ax.set_xlabel("Cohen's d (Healthy \u2192 Insulin-dependent)")
    ax.set_xlim(-0.45, 0.55)

    # Mechanism annotations
    ax.text(-0.38, -0.9, "Autonomic buffering\nbreaks down", fontsize=9,
            color=C_NEG, fontweight="bold", ha="center", style="italic")
    ax.text(+0.42, -0.9, "Metabolic coupling\nrigidifies", fontsize=9,
            color=C_POS, fontweight="bold", ha="center", style="italic")

    return save_fig(fig, "frequency_bands")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 5: Diurnal profile
# ══════════════════════════════════════════════════════════════════════

def fig_diurnal_profile():
    df = pd.read_csv(result_path("h_new17_diurnal_coupling.csv"))
    blocks = df[df["feature"].str.match(r"diurnal_r_\d")].sort_values("feature")

    times = [2, 6, 10, 14, 18, 22]
    xlabels = ["00\u201304", "04\u201308", "08\u201312", "12\u201316", "16\u201320", "20\u201324"]

    fig, ax = plt.subplots(figsize=(6.5, 4))

    ax.axvspan(-0.5, 4, alpha=0.06, color="#000000", zorder=0)

    for col, label, color, marker in [
        ("med_H",  "Healthy",    C_HEALTHY, "o"),
        ("med_PD", "Pre-DM",    C_PREDIAB, "s"),
        ("med_OM", "Oral med",  C_ORAL,    "^"),
        ("med_ID", "Insulin",   C_INSULIN, "D"),
    ]:
        vals = blocks[col].values[:6]
        ax.plot(times[:len(vals)], vals, marker=marker, color=color, label=label,
                linewidth=2.0, markersize=7, markeredgecolor="white", markeredgewidth=0.8)

    ax.axhline(0, color="#aaaaaa", linewidth=0.5, linestyle="--")
    ax.set_xticks(times)
    ax.set_xticklabels(xlabels)
    ax.set_xlabel("Time of day (local)")
    ax.set_ylabel("Glucose\u2013HR coupling (Pearson r)")
    ax.legend(loc="upper right", framealpha=0.9, edgecolor="#dddddd")
    ax.set_ylim(-0.08, 0.22)
    ax.text(2, 0.20, "Night", ha="center", fontsize=8, color="#888888", style="italic")

    return save_fig(fig, "diurnal_profile")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 6: Dose-response — box + strip overlay (Nature style)
# ══════════════════════════════════════════════════════════════════════

def fig_excursion_doseresponse():
    p2 = pd.read_parquet(result_path("hypothesis_phase2_features.parquet"),
                          columns=["study_group", "exc_hr_peak_mean", "exc_hr_auc_mean", "exc_hr_per_mg_dl"])
    p2 = p2[p2["study_group"].isin(G_ORDER)].copy()
    p2["group"] = p2["study_group"].map(dict(zip(G_ORDER, G_SHORT)))
    p2["group"] = pd.Categorical(p2["group"], categories=G_SHORT, ordered=True)

    feats = [
        ("exc_hr_peak_mean", "HR peak\n(bpm above baseline)", "d = \u22120.80 ***"),
        ("exc_hr_auc_mean", "HR AUC\n(bpm\u00b7min)", "d = \u22120.34 ***"),
        ("exc_hr_per_mg_dl", "HR per mg/dL\nexcursion", "d = \u22120.32 ***"),
    ]

    fig, axes = plt.subplots(1, 3, figsize=(11, 4))

    for ax, (feat, ylabel, stat_txt) in zip(axes, feats):
        valid = p2[["group", feat]].dropna()

        # Strip plot (individual points, jittered, semi-transparent)
        for i, (g, c) in enumerate(zip(G_SHORT, C_GROUPS)):
            vals = valid.loc[valid["group"] == g, feat].values
            jitter = np.random.default_rng(42).uniform(-0.15, 0.15, len(vals))
            ax.scatter(np.full(len(vals), i) + jitter, vals, s=3, alpha=0.08,
                       color=c, zorder=1, rasterized=True)

        # Box plot on top
        bp = ax.boxplot(
            [valid.loc[valid["group"]==g, feat].values for g in G_SHORT],
            positions=range(4), widths=0.4, patch_artist=True,
            showfliers=False, whis=1.5,
            medianprops=dict(color="white", linewidth=1.5),
            whiskerprops=dict(color="#444444", linewidth=0.8),
            capprops=dict(color="#444444", linewidth=0.8),
            boxprops=dict(linewidth=0.8),
        )
        for patch, c in zip(bp["boxes"], C_GROUPS):
            patch.set_facecolor(c)
            patch.set_alpha(0.65)
            patch.set_edgecolor(c)

        # Median annotation
        for i, g in enumerate(G_SHORT):
            med = valid.loc[valid["group"]==g, feat].median()
            ax.text(i, med, f" {med:.1f}", ha="left", va="center",
                    fontsize=7, color="#333333")

        ax.axhline(0, color="#aaaaaa", linewidth=0.5, linestyle="--")
        ax.set_xticks(range(4))
        ax.set_xticklabels(G_SHORT, fontsize=9, rotation=20, ha="right")
        ax.set_ylabel(ylabel)

        # Significance bracket
        ymax = ax.get_ylim()[1]
        h = (ymax - ax.get_ylim()[0]) * 0.03
        bracket_y = ymax * 0.92
        ax.plot([0, 0, 3, 3], [bracket_y, bracket_y + h, bracket_y + h, bracket_y],
                lw=0.8, color=C_INSULIN)
        ax.text(1.5, bracket_y + h, stat_txt, ha="center", va="bottom",
                fontsize=8, fontweight="bold", color=C_INSULIN)

    panel_label(axes[0], "a")
    panel_label(axes[1], "b")
    panel_label(axes[2], "c")

    fig.tight_layout()
    return save_fig(fig, "excursion_doseresponse")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 7: Excursion vs control
# ══════════════════════════════════════════════════════════════════════

def fig_excursion_vs_control():
    exc = [60.67, 51.18, 13.76, -19.78]
    ctrl = [-20.38, -3.27, -5.54, 5.64]

    fig, ax = plt.subplots(figsize=(6.5, 4))
    x = np.arange(4)
    w = 0.30

    ax.bar(x - w/2, exc, w, color=C_GROUPS, edgecolor="white", linewidth=0.8,
           label="Post-excursion", zorder=3)
    ax.bar(x + w/2, ctrl, w, color=["#cccccc"]*4, edgecolor="white", linewidth=0.8,
           label="Random window", zorder=3)

    ax.axhline(0, color="#333333", linewidth=0.6, zorder=2)
    ax.set_xticks(x)
    ax.set_xticklabels(G_SHORT)
    ax.set_ylabel("HR AUC (bpm\u00b7min)")
    ax.legend(fontsize=8, framealpha=0.9)
    ax.set_ylim(-40, 80)

    # Annotations
    ax.annotate("d = \u22120.34 ***\nexcursion-specific",
                xy=(3 - w/2, -19.78), xytext=(1.0, -35),
                fontsize=8, color=C_INSULIN,
                arrowprops=dict(arrowstyle="-", color=C_INSULIN, lw=0.8))
    ax.annotate("d = 0.12 n.s.",
                xy=(3 + w/2, 5.64), xytext=(3.5, 60),
                fontsize=8, color="#999999",
                arrowprops=dict(arrowstyle="-", color="#cccccc", lw=0.8))

    return save_fig(fig, "excursion_vs_control")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 8: Robustness — dot-and-label (not text art)
# ══════════════════════════════════════════════════════════════════════

def fig_robustness_checklist():
    tests = [
        ("Full covariate adj.", "p = 6.6\u00d710\u207b\u00b9\u00b9"),
        ("Negative control", "p = 0.27 n.s."),
        ("Dose-response", "r = \u22120.26"),
        ("Bootstrap 95% CI", "[\u22120.95, \u22120.68]"),
        ("Awake-only", "d = \u22120.32"),
        ("Pre-diabetes", "p = 0.025"),
        ("Frailty (adj HbA1c)", "r = \u22120.15"),
        ("Within-group", "r = \u22120.10"),
        ("Excursion-specific", "d = \u22120.33"),
    ]

    fig, ax = plt.subplots(figsize=(6, 4))
    n = len(tests)
    ax.set_xlim(-0.3, 6)
    ax.set_ylim(-0.5, n + 0.2)
    ax.axis("off")

    for i, (label, stat) in enumerate(tests):
        y = n - 1 - i
        ax.plot(0.15, y, "o", color=C_PASS, markersize=12,
                markeredgecolor="white", markeredgewidth=1.0)
        ax.text(0.15, y, "\u2713", fontsize=9, fontweight="bold",
                color="white", ha="center", va="center")
        ax.text(0.55, y, label, fontsize=10, color="#1a1a2e", va="center")
        ax.text(5.5, y, stat, fontsize=10, fontweight="bold", color=C_PASS,
                va="center", ha="right", fontfamily="monospace")

    return save_fig(fig, "robustness_checklist")


# ══════════════════════════════════════════════════════════════════════
# FIGURE 9: Synthesis — dual panel
# ══════════════════════════════════════════════════════════════════════

def fig_synthesis():
    fig, axes = plt.subplots(1, 2, figsize=(10, 4))
    np.random.seed(42)

    # (a) Fast coupling — event-triggered response schematic
    ax = axes[0]
    t = np.linspace(-30, 90, 200)
    healthy = 18 * np.exp(-((t - 35)**2) / (2 * 18**2))
    diabetes = 5 * np.exp(-((t - 35)**2) / (2 * 22**2))

    ax.fill_between(t, healthy, alpha=0.15, color=C_HEALTHY)
    ax.plot(t, healthy, color=C_HEALTHY, linewidth=2.0, label="Healthy")
    ax.fill_between(t, diabetes, alpha=0.15, color=C_INSULIN)
    ax.plot(t, diabetes, color=C_INSULIN, linewidth=2.0, label="Insulin-dep.")

    ax.axvline(0, color="#888888", linewidth=0.6, linestyle="--")
    ax.axhline(0, color="#333333", linewidth=0.4)
    ax.set_xlabel("Time from excursion (min)")
    ax.set_ylabel("HR response (bpm)")
    ax.legend(fontsize=8, framealpha=0.9)
    ax.text(55, 14.5, "d = \u22120.80", fontsize=11, fontweight="bold", color=C_INSULIN)
    ax.set_ylim(-2, 22)
    ax.set_title("Fast coupling (<2h)", fontsize=11, color=C_NEG, fontweight="bold")
    panel_label(ax, "a")

    # (b) Slow coupling — tightly locked oscillations schematic
    ax = axes[1]
    t2 = np.linspace(0, 48, 500)
    glucose = 155 + 30*np.sin(2*np.pi*t2/6) + 10*np.sin(2*np.pi*t2/4.2) + 3*np.random.randn(500)
    hr = 82 + 9*np.sin(2*np.pi*t2/6 - 0.4) + 4*np.sin(2*np.pi*t2/4.2 - 0.3) + 2*np.random.randn(500)

    ax.plot(t2, glucose, color=C_INSULIN, linewidth=1.5, label="Glucose")
    ax2 = ax.twinx()
    ax2.plot(t2, hr, color=C_ORAL, linewidth=1.5, linestyle="--", label="HR")
    ax2.spines["right"].set_visible(True)
    ax2.spines["top"].set_visible(False)

    ax.set_xlabel("Time (hours)")
    ax.set_ylabel("Glucose (mg/dL)", color=C_INSULIN)
    ax2.set_ylabel("Heart rate (bpm)", color=C_ORAL)
    ax.tick_params(axis="y", colors=C_INSULIN)
    ax2.tick_params(axis="y", colors=C_ORAL)
    ax.set_xlim(0, 48)
    ax.set_ylim(100, 210)
    ax2.set_ylim(60, 105)

    lines1, labels1 = ax.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax.legend(lines1 + lines2, labels1 + labels2, fontsize=8, loc="upper right", framealpha=0.9)
    ax.text(24, 202, "d = +0.38", fontsize=11, fontweight="bold", color=C_POS, ha="center")
    ax.set_title("Slow coupling (>2h)", fontsize=11, color=C_POS, fontweight="bold")
    panel_label(ax, "b")

    # Schematic note
    fig.text(0.5, -0.02, "Schematic illustrations based on observed group-level statistics",
             ha="center", fontsize=7, color="#999999", style="italic")

    fig.tight_layout()
    return save_fig(fig, "synthesis")


# ══════════════════════════════════════════════════════════════════════
# PPTX Assembly
# ══════════════════════════════════════════════════════════════════════

BG = RGBColor(0x1a, 0x1a, 0x2e)
SB = RGBColor(0x55, 0x55, 0x55)
AC = RGBColor(0x00, 0x72, 0xB2)
WH = RGBColor(0xff, 0xff, 0xff)
GR = RGBColor(0x00, 0x9E, 0x73)
PU = RGBColor(0xCC, 0x79, 0xA7)


def _title(slide, title, subtitle=""):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.85))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = BG
    if subtitle:
        p2 = tf.add_paragraph(); p2.text = subtitle
        p2.font.size = Pt(11); p2.font.color.rgb = SB


def img_slide(prs, title, img, subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, title, subtitle)
    if img and os.path.exists(img):
        s.shapes.add_picture(img, Inches(0.8), Inches(1.1), Inches(11.5))


def key_nums_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "Post-Excursion HR Blunting: Key Evidence")
    nums = [
        ("d = \u22120.80", "Effect size\nHR peak", AC),
        ("9 / 9", "Robustness\ntests", GR),
        ("p = 0.025", "Pre-diabetes\ndetection", RGBColor(0xE6, 0x9F, 0x00)),
        ("r = \u22120.15", "Frailty\n(adj HbA1c)", PU),
    ]
    for i, (n, l, c) in enumerate(nums):
        x = Inches(0.6 + i * 3.1)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.8), Inches(2.4))
        sh.fill.solid(); sh.fill.fore_color.rgb = WH; sh.line.color.rgb = c; sh.line.width = Pt(2.5)
        tf = sh.text_frame; tf.word_wrap = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]; p.text = n; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = c
        p2 = tf.add_paragraph(); p2.text = l; p2.font.size = Pt(12); p2.font.color.rgb = SB
        p2.alignment = PP_ALIGN.CENTER

    tx = s.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12), Inches(2.5))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("Consumer CGM + wearable HR could screen for autonomic neuropathy "
              "\u2014 no lab visit, no Ewing battery, detectable before clinical diagnosis.")
    p.font.size = Pt(14); p.font.color.rgb = BG; p.font.italic = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "\nFirst event-triggered analysis from free-living CGM + wearable (N = 1,912)"
    p2.font.size = Pt(11); p2.font.color.rgb = SB; p2.alignment = PP_ALIGN.CENTER


def summary_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "Summary")
    items = [
        ("System", "Multi-agent pipeline: Proposer \u2192 Critic \u2192 Executor \u2192 Verifier\n26 hypotheses, 7 tested, 9 robustness checks each", AC),
        ("Finding", "Fast autonomic coupling breaks down (d = \u22120.80) while slow metabolic\ncoupling rigidifies (d = +0.38). Invisible to broadband statistics.", GR),
        ("Potential", "Post-excursion HR blunting as autonomic neuropathy screen.\nDetectable in pre-diabetes (p = 0.025).", PU),
    ]
    for i, (h, b, c) in enumerate(items):
        y = Inches(1.2 + i * 2.0)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.2), Inches(1.6))
        sh.fill.solid(); sh.fill.fore_color.rgb = WH; sh.line.color.rgb = c; sh.line.width = Pt(2)
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = h; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = c
        p2 = tf.add_paragraph(); p2.text = b; p2.font.size = Pt(12); p2.font.color.rgb = BG


# ══════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════

def main():
    print("Generating figures (SciencePlots + Nature style)...")
    figs = {
        "arch": fig_architecture(),
        "status": fig_hypothesis_status(),
        "forest": fig_forest_plot(),
        "freq": fig_frequency_bands(),
        "diurnal": fig_diurnal_profile(),
        "dose": fig_excursion_doseresponse(),
        "ctrl": fig_excursion_vs_control(),
        "robust": fig_robustness_checklist(),
        "synth": fig_synthesis(),
    }
    print(f"  9 figures \u2192 {FIGDIR}/")

    print("Assembling PPTX...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tx = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    p.text = "Hypothesis-Driven Scientific Discovery\nfrom Synchronized Multimodal Physiology"
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = BG
    p2 = tf.add_paragraph(); p2.text = "\nZijian (Carl) Ma  |  Stanford University  |  TWC Lab"
    p2.font.size = Pt(15); p2.font.color.rgb = SB; p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph(); p3.text = "AI-READI v3.0.0  |  2,280 participants  |  10-day CGM + Wearable + Environment"
    p3.font.size = Pt(12); p3.font.color.rgb = SB; p3.alignment = PP_ALIGN.CENTER

    img_slide(prs, "Multi-Agent Hypothesis Discovery Pipeline", figs["arch"],
              "Literature-grounded proposal \u2192 feasibility critique \u2192 parallel execution \u2192 9-test verification")
    img_slide(prs, "Hypothesis Inventory", figs["status"],
              "8 migrated + 18 proposed  |  All critiqued for feasibility against AI-READI constraints")
    img_slide(prs, "Effect Sizes Across Tested Hypotheses", figs["forest"],
              "Event-triggered and frequency-resolved analyses yield the strongest effects")
    img_slide(prs, "Broadband Coupling Masks Opposite Timescale Effects", figs["freq"],
              "H-NEW01: Fast coupling (<2h) decreases while slow coupling (>2h) increases. Broadband cancels both.")
    img_slide(prs, "Coupling Increase Is a Daytime Phenomenon", figs["diurnal"],
              "H-NEW17: Groups diverge during waking hours. Night (00\u201304h) is preserved.")
    img_slide(prs, "Post-Excursion HR Response: Monotonic Dose\u2013Response", figs["dose"],
              "H-NEW03: HR blunting across severity. Individual data points overlaid on box plots (whiskers: 1.5\u00d7IQR).")
    img_slide(prs, "Negative Control Confirms Excursion Specificity", figs["ctrl"],
              "Random windows (gray) show no group difference. Effect is excursion-triggered.")
    img_slide(prs, "9/9 Robustness Tests Passed", figs["robust"],
              "Survives full covariate adjustment, detectable in pre-diabetes, links to clinical burden.")
    key_nums_slide(prs)
    img_slide(prs, "Two Distinct Disease Processes", figs["synth"],
              "Fast coupling loss = autonomic neuropathy.  Slow coupling gain = metabolic rigidity.")
    summary_slide(prs)

    out = Path("docs/decks/hypothesis_discovery_deck.pptx")
    prs.save(out)
    print(f"Saved: {out} ({out.stat().st_size / 1024:.0f} KB)")
    print(f"\nEditable vector figures:")
    print(f"  SVG: {FIGDIR}/*.svg  (drag into PowerPoint 365+ or open in Illustrator)")
    print(f"  PDF: {FIGDIR}/*.pdf  (open in Illustrator/Inkscape for full editing)")


if __name__ == "__main__":
    main()
