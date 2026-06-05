#!/usr/bin/env python3
"""Generate AI-READI dataset overview slide deck — ICML-style white academic aesthetic."""

import sys
sys.path.insert(0, "/oak/stanford/scg/lab_twc/mazijian/aireadi")

import io, textwrap
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.gridspec as gridspec
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from scripts.config import result_path

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
DATA = Path("/oak/stanford/scg/lab_twc/mazijian/aireadi/data")
RESULTS = Path("/oak/stanford/scg/lab_twc/mazijian/aireadi/results")
OUT = Path("/oak/stanford/scg/lab_twc/mazijian/aireadi/docs/decks/aireadi_overview.pptx")
PID = "1033"   # Example participant (healthy, 59 yr, UW, all 9 modalities)

# ---------------------------------------------------------------------------
# ICML palette: clean white, charcoal text, muted academic colours
# ---------------------------------------------------------------------------
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG          = RGBColor(0xFB, 0xFB, 0xFB)   # very slight warm white
CHARCOAL    = RGBColor(0x33, 0x33, 0x33)
GREY60      = RGBColor(0x66, 0x66, 0x66)
GREY80      = RGBColor(0x88, 0x88, 0x88)
GREYB0      = RGBColor(0xB0, 0xB0, 0xB0)
BLUE        = RGBColor(0x26, 0x6D, 0xB6)   # primary accent
RED         = RGBColor(0xC4, 0x3E, 0x3E)   # alert / highlight
TEAL        = RGBColor(0x1F, 0x8A, 0x87)
GOLD        = RGBColor(0xD4, 0x8B, 0x1A)
PURPLE      = RGBColor(0x7B, 0x52, 0xA8)

# Matplotlib
plt.rcParams.update({
    "figure.facecolor": "#FBFBFB",
    "axes.facecolor":   "#FBFBFB",
    "savefig.facecolor":"#FBFBFB",
    "text.color":       "#333333",
    "axes.labelcolor":  "#333333",
    "xtick.color":      "#555555",
    "ytick.color":      "#555555",
    "axes.edgecolor":   "#CCCCCC",
    "grid.color":       "#E8E8E8",
    "axes.grid":        False,
    "font.family":      "sans-serif",
    "font.size":        11,
})
C_BLUE  = "#266DB6"
C_RED   = "#C43E3E"
C_TEAL  = "#1F8A87"
C_GOLD  = "#D48B1A"
C_PURP  = "#7B52A8"
C_CHAR  = "#333333"
C_GREY  = "#888888"
C_LGREY = "#CCCCCC"
GRP_COLS = [C_TEAL, C_GOLD, C_RED, "#9B2D5F"]   # H, Pre, Oral, Ins

# ---------------------------------------------------------------------------
# Data loading
# ---------------------------------------------------------------------------
participants = pd.read_csv(DATA / "participants.tsv", sep="\t")
table1 = pd.read_csv(result_path("table1_by_study_group.csv"))
all_feats = pd.read_csv(result_path("all_features_by_study_group.csv"))
fm = pd.read_parquet(result_path("feature_matrix.parquet"))

GROUP_SHORT = {
    "healthy": "Healthy",
    "pre_diabetes_lifestyle_controlled": "Pre-diabetes",
    "oral_medication_and_or_non_insulin_injectable_medication_controlled": "Oral Meds",
    "insulin_dependent": "Insulin",
}
participants["group_short"] = participants["study_group"].map(GROUP_SHORT)
groups = ["Healthy", "Pre-diabetes", "Oral Meds", "Insulin"]

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def fig_to_buf(fig, dpi=220):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor(), edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


def add_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = BG


def _tf(slide, left, top, width, height, text, size=18,
        colour=CHARCOAL, bold=False, italic=False, align=PP_ALIGN.LEFT,
        font="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.color.rgb = colour
    p.font.bold = bold; p.font.italic = italic; p.font.name = font
    p.alignment = align
    return tf


def _bullets(slide, left, top, width, height, items, size=16,
             colour=GREY60, leading=Pt(7)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = colour
        p.font.name = "Calibri"; p.space_before = leading; p.space_after = Pt(2)
    return tf


def _accent_line(slide, left, top, width, colour=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(left), Inches(top), Inches(width), Pt(2.5))
    s.fill.solid(); s.fill.fore_color.rgb = colour; s.line.fill.background()


def _img(slide, buf, left, top, width):
    slide.shapes.add_picture(buf, Inches(left), Inches(top), Inches(width))


def _slide_header(slide, title, subtitle=None, accent_col=BLUE):
    add_bg(slide)
    _tf(slide, 0.7, 0.35, 11, 0.65, title, size=32, bold=True, colour=CHARCOAL)
    _accent_line(slide, 0.7, 0.9, 2.8, colour=accent_col)
    if subtitle:
        _tf(slide, 0.7, 1.05, 11, 0.4, subtitle, size=15, colour=GREY80, italic=True)

# ---------------------------------------------------------------------------
# Presentation (16:9)
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========================== SLIDE 1: TITLE =================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
_accent_line(sl, 0.9, 2.55, 5.0)
_tf(sl, 0.9, 2.7, 10, 1.0, "AI-READI", size=54, bold=True, colour=BLUE)
_tf(sl, 0.9, 3.65, 11, 0.75,
    "Artificial Intelligence Ready and Equitable Atlas\nfor Diabetes Insights  \u2014  v3.0.0",
    size=22, colour=CHARCOAL)
_tf(sl, 0.9, 4.7, 9, 0.5,
    "Dataset Overview & Analysis Infrastructure", size=18, colour=GREY60)
_tf(sl, 0.9, 5.7, 8, 0.4,
    "Zijian (Carl) Ma  \u00b7  Stanford TWC Lab  \u00b7  April 2026", size=15, colour=GREY80)
_tf(sl, 0.9, 6.15, 8, 0.35,
    "NIH Bridge2AI  \u00b7  1OT2OD032644  \u00b7  DOI: 10.60775/fairhub.3", size=12, colour=GREYB0)

# Stat boxes (right)
stats = [("2,280", "Participants"), ("9", "Modalities"),
         ("3.82 TB", "Dataset Size"), ("356,343", "Files")]
for i, (val, label) in enumerate(stats):
    bx, by = 9.0, 2.2 + i * 1.15
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(bx), Inches(by), Inches(3.4), Inches(0.9))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); s.line.width = Pt(0.75)
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = val
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = BLUE
    p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = label
    p2.font.size = Pt(13); p2.font.color.rgb = GREY80
    p2.font.name = "Calibri"; p2.alignment = PP_ALIGN.CENTER

# ========================== SLIDE 2: STUDY DESIGN ==========================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Study Design")
_bullets(sl, 0.7, 1.35, 5.5, 4.0, [
    "\u2022  NIH Bridge2AI Common Fund (PI: Aaron Lee, UW)",
    "\u2022  Salutogenesis: pathway from T2DM back to health",
    "\u2022  Cross-sectional: single visit + ~10-day monitoring",
    "\u2022  3 sites: UW (798)  \u00b7  UAB (800)  \u00b7  UCSD (682)",
    "\u2022  Balanced across sex, race/ethnicity, T2DM stage",
    "\u2022  Public release: sex/race/meds redacted",
    "\u2022  ML splits: Train 1,576 / Val 352 / Test 352",
    "\u2022  Collection: Jul 2023 \u2013 May 2025",
], size=16, colour=GREY60)

fig, ax = plt.subplots(figsize=(5.2, 4.0))
counts = [776, 560, 686, 258]
bars = ax.barh(groups[::-1], counts[::-1], color=GRP_COLS[::-1],
               height=0.55, edgecolor="white", linewidth=0.5)
for bar, cnt in zip(bars, counts[::-1]):
    ax.text(bar.get_width() + 12, bar.get_y() + bar.get_height()/2,
            f"  {cnt}  ({cnt/2280*100:.0f}%)", va="center", fontsize=12, color=C_CHAR)
ax.set_xlim(0, 950); ax.set_xlabel("Participants", fontsize=12)
ax.set_title("T2DM Severity Groups (n = 2,280)", fontsize=14, pad=10, color=C_CHAR)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
fig.tight_layout()
_img(sl, fig_to_buf(fig), 6.8, 1.3, 5.8)

# ========================== SLIDE 3: DEMOGRAPHICS ==========================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Participant Demographics",
              "Mean age 60.9 \u00b1 11.2 yr  |  Range 40\u201394  |  Median 61  |  IQR 52\u201369")

fig, axes = plt.subplots(1, 2, figsize=(11.5, 4.0))
ax = axes[0]
for g, col in zip(groups, GRP_COLS):
    subset = participants.loc[participants["group_short"] == g, "age"]
    ax.hist(subset, bins=np.arange(38, 98, 3), alpha=0.55, color=col,
            label=g, edgecolor="white", linewidth=0.4)
ax.set_xlabel("Age (years)"); ax.set_ylabel("Count")
ax.set_title("Age Distribution by Group", fontsize=13, pad=8)
ax.legend(fontsize=9, frameon=True, fancybox=False, edgecolor=C_LGREY)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

ax = axes[1]
sg = participants.groupby(["clinical_site","group_short"]).size().unstack(fill_value=0)[groups]
sg.plot.bar(stacked=True, ax=ax, color=GRP_COLS, edgecolor="white", linewidth=0.4, width=0.55)
ax.set_xlabel(""); ax.set_ylabel("Participants")
ax.set_title("Site \u00d7 Group Distribution", fontsize=13, pad=8)
ax.legend(fontsize=9, frameon=True, fancybox=False, edgecolor=C_LGREY)
ax.set_xticklabels(ax.get_xticklabels(), rotation=0)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.7, 1.6, 11.8)

# ========================== SLIDE 4: 9 MODALITIES ==========================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Nine Modalities",
              "Clinical Dataset Structure (CDS) v0.1.1  |  DICOM \u00b7 WFDB \u00b7 OMOP CDM \u00b7 Open mHealth")

mods = [
    ("Clinical Data",    "OMOP CDM",    "6 CSV tables \u2014 labs, vitals, conditions, procedures",  "2,280", 100.0, C_BLUE),
    ("Cardiac ECG",      "WFDB",        "Philips TC30 \u2014 12-lead, 500 Hz, 5,500 samples (11 s)", "2,251", 98.7, C_BLUE),
    ("Retinal OCT",      "DICOM",       "4 devices \u2014 Spectralis, Maestro2, Triton, Cirrus",    "2,266", 99.4, C_PURP),
    ("Retinal OCTA",     "DICOM",       "4 devices \u2014 enface + flow cube + segmentation",       "2,264", 99.3, C_PURP),
    ("Retinal Photography","DICOM",     "6 devices \u2014 CFP, FAF, IR",                            "2,275", 99.8, C_PURP),
    ("Retinal FLIO",     "DICOM",       "Heidelberg \u2014 256\u00d7256, 1,024 temporal frames",    "1,847", 81.0, C_PURP),
    ("Wearable",         "Open mHealth","Garmin Vivosmart 5 \u2014 HR, SpO\u2082, RR, stress, sleep, activity","2,184", 95.8, C_TEAL),
    ("CGM",              "Open mHealth","Dexcom G6 \u2014 5-min glucose, ~10 days",                 "2,245", 98.5, C_TEAL),
    ("Environment",      "CSV",         "LeeLab Anura \u2014 light, PM, T/RH, VOC/NOx, 5-sec",     "2,231", 97.9, C_TEAL),
]

fig, ax = plt.subplots(figsize=(12, 5.2))
ax.set_xlim(-0.3, 13); ax.set_ylim(-0.5, len(mods) + 0.5); ax.axis("off")
hx = [0.3, 2.6, 4.8, 10.5, 11.4]
for hdr, x in zip(["Modality","Standard","Description","N","Coverage"], hx):
    ax.text(x, len(mods) + 0.1, hdr, fontsize=11, fontweight="bold", color=C_CHAR, va="bottom")
ax.axhline(y=len(mods) - 0.2, color=C_LGREY, linewidth=0.8, xmin=0.02, xmax=0.99)

for i, (name, std, desc, n, pct, col) in enumerate(mods):
    y = len(mods) - 1 - i
    ax.plot(0.1, y, "o", color=col, markersize=7)
    ax.text(hx[0], y, name, fontsize=11, color=C_CHAR, va="center", fontweight="bold")
    ax.text(hx[1], y, std, fontsize=10, color=C_GREY, va="center")
    ax.text(hx[2], y, desc, fontsize=10, color=C_GREY, va="center")
    ax.text(hx[3], y, n, fontsize=11, color=C_CHAR, va="center", ha="center")
    bar_w = pct / 100 * 1.3
    ax.barh(y, bar_w, left=hx[4], height=0.35, color=col, alpha=0.55, edgecolor="none")
    ax.text(hx[4] + bar_w + 0.08, y, f"{pct:.0f}%", fontsize=9, color=C_GREY, va="center")
    if i < len(mods) - 1:
        ax.axhline(y=y - 0.45, color="#EEEEEE", linewidth=0.5, xmin=0.02, xmax=0.99)

fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.4, 1.5, 12.5)

# ========================== SLIDE 5: TEMPORAL ALIGNMENT ====================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Temporal Alignment: The ~10-Day Overlap Window",
              "Same-day clinical snapshot anchored to concurrent wearable + CGM + environmental time-series")

fig, ax = plt.subplots(figsize=(12, 4.5))
ax.set_xlim(-1.5, 14); ax.set_ylim(-0.8, 7.5); ax.axis("off")

# Timeline
ax.axhline(y=3.5, color=C_LGREY, linewidth=1.5, xmin=0.05, xmax=0.88)
for d in range(0, 12):
    ax.plot(d + 0.5, 3.5, "|", color=C_LGREY, markersize=8)
    ax.text(d + 0.5, 3.15, f"Day {d}", fontsize=8, color="#AAAAAA", ha="center")

# Visit day line
ax.axvline(x=0.5, color=C_RED, linewidth=1.5, linestyle="--", ymin=0.12, ymax=0.88)
ax.text(0.5, 7.1, "Clinic Visit Day", fontsize=13, fontweight="bold", color=C_RED, ha="center")

# Single-timepoint modalities
for label, y, col in [("Labs / Vitals / OMOP CDM", 6.5, C_BLUE),
                       ("12-lead ECG (11 sec)", 6.0, C_BLUE),
                       ("OCT / OCTA / Photography / FLIO", 5.5, C_PURP)]:
    ax.plot(0.5, y, "s", color=col, markersize=8)
    ax.text(0.85, y, label, fontsize=12, color=col, va="center")

# Continuous modalities
for label, y, col in [("Garmin Wearable (HR, SpO\u2082, sleep, stress, activity)", 2.5, C_TEAL),
                       ("Dexcom G6 CGM (glucose every 5 min)", 1.7, C_GOLD),
                       ("Anura Environmental Sensor (light, PM, T/RH, VOC)", 0.9, C_BLUE)]:
    ax.barh(y, 10.0, left=0.3, height=0.45, color=col, alpha=0.35, edgecolor=col, linewidth=0.8)
    ax.text(10.6, y, label, fontsize=11, color=col, va="center")

ax.annotate("", xy=(0.3, 0.15), xytext=(10.3, 0.15),
            arrowprops=dict(arrowstyle="<->", color=C_CHAR, lw=1.2))
ax.text(5.3, -0.2, "~8\u201310 day triple overlap", fontsize=13, color=C_CHAR,
        ha="center", fontweight="bold")
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.3, 1.4, 12.5)

_tf(sl, 0.7, 6.1, 11.5, 0.8,
    "This dense cross-modal snapshot is rare: concurrent wearable + CGM + environment\n"
    "time-series anchored to same-day multi-device retinal imaging + ECG + 105 clinical labs.\n"
    "Enables contrastive and masked pretraining of shared latent representations.",
    size=14, colour=GREY80, italic=True)

# ========================== SLIDE 6: CLINICAL DATA =========================
print("  Generating clinical slide...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Modality 1: Clinical Data (OMOP CDM)", accent_col=BLUE)

_bullets(sl, 0.7, 1.3, 5.8, 3.5, [
    "\u2022  6 standardised CSV tables (OMOP CDM v5.4)",
    "\u2022  person, visit, measurement, observation,",
    "   condition_occurrence, procedure_occurrence",
    "\u2022  105 measurement concepts (labs + vitals)",
    "\u2022  244 observation concepts (questionnaires)",
    "\u2022  30 condition flags (ICD-10 mapped)",
    "\u2022  Feature matrix: 2,280 \u00d7 125 (Parquet cached)",
], size=15, colour=GREY60)

# Show mini table of example values
from scripts.loaders.clinical import get_measurements
from scripts.utils.concepts import get_concept_label
m = get_measurements(PID)

fig, ax = plt.subplots(figsize=(5.5, 4.5))
ax.axis("off"); ax.set_xlim(0, 10); ax.set_ylim(-0.5, 11)
ax.text(0.2, 10.5, f"Example: Participant {PID} (Healthy, 59 yr)", fontsize=13,
        fontweight="bold", color=C_CHAR)
ax.axhline(y=10.1, color=C_LGREY, linewidth=0.8, xmin=0.02, xmax=0.98)

# Pick representative labs
example_concepts = [
    (3004410, "HbA1c"), (3000963, "Glucose"), (3022192, "Triglycerides"),
    (3027114, "HDL"), (3028437, "LDL"), (3020891, "BMI"),
    (3004249, "SBP"), (3012888, "DBP"), (3027018, "Heart rate"),
]
show_items = []
for cid, short_name in example_concepts:
    row = m[m["measurement_concept_id"] == cid]
    if len(row) > 0:
        val = row.iloc[0]["value_as_number"]
        show_items.append((short_name, f"{val:.1f}" if pd.notna(val) else "N/A"))

for i, (name, val) in enumerate(show_items):
    y = 9.5 - i * 1.05
    ax.text(0.4, y, name, fontsize=11, color=C_CHAR, va="center")
    ax.text(6.5, y, val, fontsize=11, color=C_BLUE, va="center", fontweight="bold", ha="right")
    if i < len(show_items) - 1:
        ax.axhline(y=y - 0.5, color="#EEEEEE", linewidth=0.4, xmin=0.04, xmax=0.96)

fig.tight_layout()
_img(sl, fig_to_buf(fig), 7.0, 1.2, 5.8)

# ========================== SLIDE 7: CARDIAC ECG ===========================
print("  Generating ECG slide...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Modality 2: Cardiac ECG", accent_col=BLUE)

_bullets(sl, 0.7, 1.3, 4.5, 2.5, [
    "\u2022  Philips PageWriter TC30",
    "\u2022  12 leads \u00d7 500 Hz \u00d7 5,500 samples (11 sec)",
    "\u2022  WFDB format (.hea + .dat)",
    "\u2022  2,251 participants (98.7%)",
    "\u2022  Device-reported: Rate, PR, QRS, QT, QTc",
], size=15, colour=GREY60)

from scripts.loaders.ecg import load_ecg, ECG_LEAD_NAMES
recs = load_ecg(PID)
sig, meta = recs[0]

fig, axes = plt.subplots(6, 2, figsize=(11, 5.0), sharex=True)
t = np.arange(sig.shape[0]) / 500.0
for i in range(12):
    ax = axes[i // 2, i % 2]
    ax.plot(t, sig[:, i], color=C_BLUE, linewidth=0.4, alpha=0.85)
    ax.set_ylabel(ECG_LEAD_NAMES[i], fontsize=8, rotation=0, labelpad=22, va="center")
    ax.set_ylim(sig[:, i].mean() - 1.5, sig[:, i].mean() + 1.5)
    ax.tick_params(axis="both", labelsize=7)
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
axes[5, 0].set_xlabel("Time (sec)", fontsize=9)
axes[5, 1].set_xlabel("Time (sec)", fontsize=9)
fig.suptitle(f"Participant {PID}  \u2014  HR {meta.get('Rate', '?')} bpm, QTc {meta.get('QTc', '?')} ms",
             fontsize=13, color=C_CHAR, y=1.01)
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.5, 2.2, 12.3)

# ========================== SLIDE 8: CGM ===================================
print("  Generating CGM slide...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Modality 3: Continuous Glucose Monitoring (CGM)", accent_col=TEAL)

from scripts.loaders.cgm import load_cgm, compute_cgm_metrics
df_cgm, hdr = load_cgm(PID)
metrics = compute_cgm_metrics(df_cgm)

fig = plt.figure(figsize=(12, 4.8))
gs = gridspec.GridSpec(1, 4, width_ratios=[3, 1, 0.1, 1.2])

# Glucose trace
ax = fig.add_subplot(gs[0])
ax.plot(df_cgm.index, df_cgm["glucose_mg_dl"], color=C_TEAL, linewidth=0.6, alpha=0.8)
ax.axhspan(70, 180, alpha=0.08, color=C_TEAL, label="Target range (70\u2013180)")
ax.axhline(180, color=C_GOLD, linewidth=0.5, linestyle="--", alpha=0.5)
ax.axhline(70, color=C_RED, linewidth=0.5, linestyle="--", alpha=0.5)
ax.set_ylabel("Glucose (mg/dL)", fontsize=11)
ax.set_xlabel("Date", fontsize=11)
ax.set_title(f"Participant {PID} \u2014 Dexcom G6, {metrics['duration_days']:.0f} days, "
             f"{metrics['n_readings']} readings", fontsize=12, pad=8)
ax.legend(fontsize=9, loc="upper right", frameon=True, edgecolor=C_LGREY)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
ax.tick_params(axis="x", rotation=30, labelsize=8)

# TIR stacked bar
ax2 = fig.add_subplot(gs[1])
tir_vals = [metrics.get("tbr_vlow", 0), metrics.get("tbr_low", 0),
            metrics["tir"], metrics.get("tar_high", 0), metrics.get("tar_vhigh", 0)]
tir_labels = ["<54", "54\u201369", "70\u2013180", "181\u2013250", ">250"]
tir_colors = ["#8B0000", C_RED, C_TEAL, C_GOLD, "#D4551A"]
bottom = 0
for val, lab, col in zip(tir_vals, tir_labels, tir_colors):
    ax2.bar(0, val, bottom=bottom, color=col, width=0.5, edgecolor="white", linewidth=0.5)
    if val > 0.03:
        ax2.text(0, bottom + val/2, f"{val*100:.0f}%", ha="center", va="center",
                fontsize=9, color="white", fontweight="bold")
    bottom += val
ax2.set_xlim(-0.5, 0.5); ax2.set_ylim(0, 1)
ax2.set_title("TIR", fontsize=12, pad=8)
ax2.set_xticks([]); ax2.set_ylabel("Fraction", fontsize=10)
ax2.spines["top"].set_visible(False); ax2.spines["right"].set_visible(False)

# Metrics table
ax3 = fig.add_subplot(gs[3])
ax3.axis("off"); ax3.set_xlim(0, 4); ax3.set_ylim(-0.5, 8)
ax3.text(0.1, 7.5, "Summary", fontsize=12, fontweight="bold", color=C_CHAR)
met_items = [
    ("Mean glucose", f"{metrics['mean_glucose']:.1f} mg/dL"),
    ("SD", f"{metrics['std_glucose']:.1f} mg/dL"),
    ("CV", f"{metrics['cv']:.1%}"),
    ("TIR", f"{metrics['tir']:.1%}"),
    ("GMI", f"{metrics['gmi']:.1f}%"),
    ("MAGE", f"{metrics['mage']:.1f} mg/dL"),
]
for i, (k, v) in enumerate(met_items):
    y = 6.5 - i * 1.05
    ax3.text(0.1, y, k, fontsize=10, color=C_GREY)
    ax3.text(3.8, y, v, fontsize=10, color=C_CHAR, ha="right", fontweight="bold")

fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.3, 1.4, 12.5)

_tf(sl, 0.7, 6.5, 11, 0.4,
    "Dexcom G6  \u00b7  5-min sampling  \u00b7  Open mHealth blood-glucose v3.0 JSON  "
    "\u00b7  2,245 participants (98.5%)  \u00b7  ~9\u201312 days per person",
    size=13, colour=GREY80, italic=True)

# ========================== SLIDE 9: WEARABLE ==============================
print("  Generating wearable slide...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Modality 4: Wearable Activity Monitor", accent_col=TEAL)

from scripts.loaders.wearable import load_wearable
w = load_wearable(PID)
hr = w["heart_rate"]
spo2 = w["oxygen_saturation"]
stress = w["stress"]

fig, axes = plt.subplots(3, 1, figsize=(12, 5.0), sharex=True)

# HR
ax = axes[0]
ax.plot(hr.index, hr["value"], color=C_RED, linewidth=0.3, alpha=0.6)
ax.set_ylabel("Heart Rate\n(bpm)", fontsize=10)
ax.set_title(f"Participant {PID} \u2014 Garmin Vivosmart 5 ({len(hr):,} HR readings over ~20 days)",
             fontsize=12, pad=6)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

# SpO2
ax = axes[1]
ax.plot(spo2.index, spo2["value"], color=C_BLUE, linewidth=0.3, alpha=0.6)
ax.set_ylabel("SpO\u2082 (%)", fontsize=10)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

# Stress
ax = axes[2]
ax.plot(stress.index, stress["value"], color=C_PURP, linewidth=0.3, alpha=0.6)
ax.set_ylabel("Stress\n(0\u2013100)", fontsize=10)
ax.set_xlabel("Date (UTC)", fontsize=10)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

for a in axes:
    a.tick_params(axis="x", rotation=20, labelsize=8)
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.3, 1.4, 12.5)

_tf(sl, 0.7, 6.6, 11, 0.4,
    "7 sub-modalities: heart_rate, SpO\u2082, respiratory_rate, stress, sleep, "
    "physical_activity, calories  \u00b7  Open mHealth JSON  \u00b7  2,184 participants (95.8%)",
    size=13, colour=GREY80, italic=True)

# ========================== SLIDE 10: ENVIRONMENT ==========================
print("  Generating environment slide...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Modality 5: Environmental Sensor", accent_col=BLUE)

from scripts.loaders.environment import load_environment
env_df, env_meta = load_environment(PID)

fig, axes = plt.subplots(4, 1, figsize=(12, 5.5), sharex=True)

ax = axes[0]
ax.plot(env_df.index, env_df["lch0"], color=C_GOLD, linewidth=0.3, alpha=0.7, label="lch0 (vis)")
ax.plot(env_df.index, env_df["lch10"], color=C_BLUE, linewidth=0.3, alpha=0.7, label="lch10 (blue)")
ax.set_ylabel("Light (0\u20131)", fontsize=9); ax.legend(fontsize=8, ncol=2, loc="upper right")
ax.set_title(f"Participant {PID} \u2014 LeeLab Anura, 5-sec sampling, {len(env_df):,} records",
             fontsize=12, pad=6)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

ax = axes[1]
ax.plot(env_df.index, env_df["pm2.5"], color=C_RED, linewidth=0.3, alpha=0.6)
ax.set_ylabel("PM2.5\n(\u03bcg/m\u00b3)", fontsize=9)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

ax = axes[2]
ax.plot(env_df.index, env_df["temp"], color=C_TEAL, linewidth=0.3, alpha=0.6)
ax.set_ylabel("Temp (\u00b0C)", fontsize=9)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

ax = axes[3]
ax.plot(env_df.index, env_df["hum"], color=C_PURP, linewidth=0.3, alpha=0.6)
ax.set_ylabel("Humidity (%)", fontsize=9)
ax.set_xlabel("Date (UTC)", fontsize=10)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

for a in axes:
    a.tick_params(axis="x", rotation=20, labelsize=8)
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.3, 1.3, 12.5)

_tf(sl, 0.7, 6.9, 11, 0.35,
    "22 channels: 10 spectral light (lch0\u2013lch11), PM1/2.5/4/10, humidity, temperature, "
    "VOC, NOx, screen, flicker  \u00b7  2,231 participants (97.9%)",
    size=13, colour=GREY80, italic=True)

# ========================== SLIDE 11: RETINAL IMAGING ======================
print("  Generating retinal slide (loading DICOMs)...")
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Modalities 6\u20139: Retinal Imaging", accent_col=PURPLE)

import pydicom
from scripts.loaders.retinal import get_oct_files, get_photography_files

photo = get_photography_files(PID)
oct_df = get_oct_files(PID)

# Load CFP (color fundus)
cfp_fp = str(DATA / photo[photo["filepath"].str.contains("eidon_mosaic_cfp_r")].iloc[0]["filepath"].lstrip("/"))
cfp_px = pydicom.dcmread(cfp_fp).pixel_array

# Load FAF
faf_fp = str(DATA / photo[photo["filepath"].str.contains("faf")].iloc[0]["filepath"].lstrip("/"))
faf_px = pydicom.dcmread(faf_fp).pixel_array

# Load OCT (central B-scan from volume)
oct_fp = str(DATA / oct_df[oct_df["filepath"].str.contains("cirrus_mac_oct")].iloc[0]["filepath"].lstrip("/"))
oct_px = pydicom.dcmread(oct_fp).pixel_array
mid_slice = oct_px.shape[0] // 2

# Load IR
ir_fp = str(DATA / photo[photo["filepath"].str.contains("cirrus_mac_oct_ir_r")].iloc[0]["filepath"].lstrip("/"))
ir_px = pydicom.dcmread(ir_fp).pixel_array

fig, axes = plt.subplots(1, 4, figsize=(13, 3.5))

ax = axes[0]
ax.imshow(cfp_px)
ax.set_title("Color Fundus (CFP)", fontsize=11, pad=6, color=C_CHAR)
ax.axis("off")

ax = axes[1]
ax.imshow(faf_px)
ax.set_title("Autofluorescence (FAF)", fontsize=11, pad=6, color=C_CHAR)
ax.axis("off")

ax = axes[2]
if ir_px.ndim == 3 and ir_px.shape[-1] == 3:
    ax.imshow(ir_px)
else:
    ax.imshow(ir_px, cmap="gray")
ax.set_title("Infrared (IR)", fontsize=11, pad=6, color=C_CHAR)
ax.axis("off")

ax = axes[3]
ax.imshow(oct_px[mid_slice], cmap="gray")
ax.set_title(f"OCT B-scan (slice {mid_slice}/{oct_px.shape[0]})", fontsize=11, pad=6, color=C_CHAR)
ax.axis("off")

fig.suptitle(f"Participant {PID} (right eye)", fontsize=13, color=C_CHAR, y=1.02)
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.2, 1.5, 12.8)

# Sub-modality summary
_bullets(sl, 0.7, 5.2, 6, 2.0, [
    "\u2022  OCT: structural cross-sections, 4 devices, 2,266 (99.4%)",
    "\u2022  OCTA: vascular angiography + flow cubes, 2,264 (99.3%)",
    "\u2022  Photography: CFP/FAF/IR, 6 devices, 2,275 (99.8%)",
    "\u2022  FLIO: fluorescence lifetime, 256\u00d7256\u00d71,024, 1,847 (81.0%)",
], size=14, colour=GREY60)

_bullets(sl, 7.0, 5.2, 5.5, 2.0, [
    "\u2022  Multi-device: Spectralis, Maestro2, Triton, Cirrus,",
    "   Eidon, Optomed Aurora",
    "\u2022  Both eyes (OD + OS), multiple scan protocols",
    "\u2022  DICOM format with rich metadata headers",
], size=14, colour=GREY60)

# ========================== SLIDE 12: FEATURE MATRIX =======================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Clinical Feature Matrix: 2,280 \u00d7 125",
              "Wide-format pivot of OMOP CDM long tables \u2192 one row per participant")

blocks = [
    ("Labs (38)", "HbA1c, glucose, insulin, C-peptide, lipid panel, CBC,\nmetabolic panel, hepatic, renal, cardiac markers", C_BLUE),
    ("Vitals (9)", "SBP, DBP, heart rate, height, weight, BMI,\nwaist, hip, waist-to-hip ratio", C_TEAL),
    ("Vision (22)", "VA + contrast sensitivity (OD/OS, photopic/mesopic),\nautorefraction (sphere/cylinder/axis OD/OS)", C_PURP),
    ("Cognition (16)", "MoCA total + 15 subscores (memory, trails, clock,\ncube, naming, attention, abstraction, orientation)", C_GOLD),
    ("Conditions (30)", "Boolean flags: T2DM, hypertension, dyslipidemia,\ncardiovascular, arthritis, depression, neuropathy, ...", C_RED),
    ("Other (5)", "CES-D-10 depression, monofilament neuropathy,\never smoked, ever alcohol, demographics", C_GREY),
]

fig, ax = plt.subplots(figsize=(5.5, 5.0))
ax.set_xlim(0, 7); ax.set_ylim(-0.5, len(blocks) + 0.3); ax.axis("off")
for i, (title, desc, col) in enumerate(blocks):
    y = len(blocks) - 1 - i
    rect = mpatches.FancyBboxPatch((0.05, y - 0.35), 6.9, 0.85,
                                    boxstyle="round,pad=0.08",
                                    facecolor=col, alpha=0.07,
                                    edgecolor=col, linewidth=1.0)
    ax.add_patch(rect)
    ax.text(0.25, y + 0.2, title, fontsize=12, fontweight="bold", color=col, va="center")
    ax.text(0.25, y - 0.13, desc, fontsize=9, color=C_GREY, va="center")
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.3, 1.5, 5.8)

# Missingness
fig2, ax2 = plt.subplots(figsize=(5.5, 5.0))
miss_pct = (fm.isnull().sum() / len(fm) * 100).sort_values(ascending=True).tail(15)
ax2.barh(range(len(miss_pct)), miss_pct.values, color=C_RED, alpha=0.5, edgecolor="none")
ax2.set_yticks(range(len(miss_pct)))
ax2.set_yticklabels(miss_pct.index, fontsize=10)
ax2.set_xlabel("% Missing", fontsize=11)
ax2.set_title("Missing Data (Top 15 Features)", fontsize=13, pad=8)
ax2.spines["top"].set_visible(False); ax2.spines["right"].set_visible(False)
for i, v in enumerate(miss_pct.values):
    ax2.text(v + 0.08, i, f"{v:.1f}%", va="center", fontsize=9, color=C_GREY)
fig2.tight_layout()
_img(sl, fig_to_buf(fig2), 6.8, 1.5, 5.8)

_tf(sl, 0.7, 6.7, 11, 0.35,
    "Max missingness 3.0% (HbA1c)  \u00b7  47 participants missing all blood labs (logistics)  "
    "\u00b7  Structured non-random missingness patterns",
    size=13, colour=GREY80, italic=True)

# ========================== SLIDE 13: KEY FINDINGS =========================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Cohort Comparison: Key Features Across T2DM Groups",
              "Kruskal-Wallis tests with Benjamini-Hochberg FDR correction (q < 0.05)")

key_features = [
    "HbA1c (%)", "Fasting glucose (mg/dL)", "Total cholesterol (mg/dL)",
    "HDL (mg/dL)", "LDL (mg/dL)", "Triglycerides (mg/dL)",
    "CRP (mg/L)", "Troponin T (ng/mL)", "Creatinine (mg/dL)",
    "BMI (kg/m\u00b2)", "Waist circumference (cm)", "Heart rate (bpm)",
    "MoCA total score", "CES-D-10 total score",
]
t1_sel = table1[table1["feature"].isin(key_features)].copy()
t1_sel = t1_sel.set_index("feature").loc[key_features].reset_index()

fig, ax = plt.subplots(figsize=(12.5, 5.5))
ax.axis("off")
col_x = [0.0, 3.7, 5.7, 7.7, 9.7, 11.5]
headers = ["Feature", "Healthy (776)", "Pre-diab (560)", "Oral Meds (686)", "Insulin (258)", "p-value"]

for hdr, x in zip(headers, col_x):
    ax.text(x, len(t1_sel) + 0.5, hdr, fontsize=10, fontweight="bold", color=C_CHAR, va="center")
ax.axhline(y=len(t1_sel) + 0.12, color=C_LGREY, linewidth=0.8, xmin=0, xmax=1)

for i, row in t1_sel.iterrows():
    y = len(t1_sel) - 0.5 - i * 0.88
    ax.text(col_x[0], y, row["feature"], fontsize=10, color=C_CHAR, va="center")
    for j, grp in enumerate(["Healthy","Pre-diabetes","Oral/Non-insulin Meds","Insulin-dependent"]):
        m_val = row[f"mean_{grp}"]; s_val = row[f"std_{grp}"]
        ax.text(col_x[j+1], y, f"{m_val:.1f} \u00b1 {s_val:.1f}",
                fontsize=9.5, color=C_GREY, va="center")
    pv = row["pvalue"]
    if pv < 1e-10: pv_txt, pv_col = f"{pv:.0e}", C_RED
    elif pv < 0.001: pv_txt, pv_col = f"{pv:.0e}", C_GOLD
    elif pv < 0.05: pv_txt, pv_col = f"{pv:.4f}", C_CHAR
    else: pv_txt, pv_col = "ns", C_LGREY
    ax.text(col_x[5], y, pv_txt, fontsize=9.5, color=pv_col, va="center")
    if i < len(t1_sel) - 1:
        ax.axhline(y=y - 0.42, color="#F0F0F0", linewidth=0.4, xmin=0, xmax=1)

ax.set_xlim(-0.2, 12.8); ax.set_ylim(-1, len(t1_sel) + 1.2)
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.2, 1.4, 12.8)

# ========================== SLIDE 14: EFFECT SIZES =========================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Feature Discrimination: Effect Sizes & Significance",
              "79 of 89 numeric features significant after FDR correction")

fig, axes = plt.subplots(1, 2, figsize=(12.5, 5.0), gridspec_kw={"width_ratios": [1.3, 1]})

ax = axes[0]
top20 = all_feats.sort_values("effect_size", ascending=False).head(20)
colours = [C_BLUE if es > 0.06 else C_TEAL if es > 0.03 else C_GREY for es in top20["effect_size"]]
ax.barh(range(len(top20)), top20["effect_size"].values[::-1],
        color=colours[::-1], edgecolor="white", linewidth=0.3, height=0.6)
ax.set_yticks(range(len(top20)))
ax.set_yticklabels(top20["feature"].values[::-1], fontsize=9.5)
ax.set_xlabel("Effect Size (\u03b7\u00b2 / \u03b5\u00b2)", fontsize=11)
ax.set_title("Top 20 Features", fontsize=13, pad=8)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)

ax = axes[1]
sig_mask = all_feats["significant"] == True
ax.scatter(all_feats.loc[~sig_mask, "effect_size"],
           -np.log10(all_feats.loc[~sig_mask, "pvalue"].clip(lower=1e-300)),
           c=C_LGREY, s=25, alpha=0.6, edgecolors="none", label="Not sig.", zorder=2)
ax.scatter(all_feats.loc[sig_mask, "effect_size"],
           -np.log10(all_feats.loc[sig_mask, "pvalue"].clip(lower=1e-300)),
           c=C_BLUE, s=30, alpha=0.7, edgecolors="none", label="FDR < 0.05", zorder=3)
for _, row in all_feats.head(5).iterrows():
    ax.annotate(row["feature"],
                (row["effect_size"], -np.log10(max(row["pvalue"], 1e-300))),
                fontsize=8.5, color=C_RED,
                xytext=(5, 5), textcoords="offset points")
ax.set_xlabel("Effect Size", fontsize=11)
ax.set_ylabel("-log\u2081\u2080(p)", fontsize=11)
ax.set_title("Effect Size vs. Significance", fontsize=13, pad=8)
ax.legend(fontsize=9, frameon=True, edgecolor=C_LGREY)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.2, 1.4, 12.8)

# ========================== SLIDE 15: LAB TRAJECTORIES =====================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Lab Values Across T2DM Severity",
              "Mean \u00b1 SD per group for key clinical biomarkers")

key_labs = ["hba1c","glucose","total_cholesterol","hdl","triglycerides","creatinine"]
titles_ = ["HbA1c (%)","Fasting Glucose\n(mg/dL)","Total Cholesterol\n(mg/dL)",
           "HDL (mg/dL)","Triglycerides\n(mg/dL)","Creatinine (mg/dL)"]

fig, axes = plt.subplots(2, 3, figsize=(12, 5.2))
for idx, (lab, title) in enumerate(zip(key_labs, titles_)):
    ax = axes[idx // 3][idx % 3]
    t1r = table1[table1["column"] == lab]
    if t1r.empty: continue
    t1r = t1r.iloc[0]
    means = [t1r["mean_Healthy"], t1r["mean_Pre-diabetes"],
             t1r["mean_Oral/Non-insulin Meds"], t1r["mean_Insulin-dependent"]]
    stds  = [t1r["std_Healthy"], t1r["std_Pre-diabetes"],
             t1r["std_Oral/Non-insulin Meds"], t1r["std_Insulin-dependent"]]
    x = np.arange(4)
    ax.bar(x, means, yerr=stds, color=GRP_COLS, width=0.55,
           capsize=3, error_kw={"ecolor": C_LGREY, "capthick": 0.8},
           edgecolor="white", linewidth=0.5)
    ax.set_xticks(x); ax.set_xticklabels(["H","Pre","Oral","Ins"], fontsize=10)
    ax.set_title(title, fontsize=11, pad=6)
    ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
fig.tight_layout(h_pad=2.5)
_img(sl, fig_to_buf(fig), 0.5, 1.4, 12.0)

_tf(sl, 0.7, 6.8, 11, 0.35,
    "H = Healthy  |  Pre = Pre-diabetes  |  Oral = Oral/non-insulin meds  |  "
    "Ins = Insulin-dependent  |  Error bars = 1 SD",
    size=13, colour=GREY80, italic=True)

# ========================== SLIDE 16: WEARABLE FEATURES ====================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Derived Features: Wearable & CGM",
              "Timezone-aware processing  \u00b7  Robust to sensor artifacts  \u00b7  Cached Parquet output")

fig, axes = plt.subplots(1, 2, figsize=(12, 5.5))

for ax, title, items, col in [
    (axes[0], "Wearable (Garmin Vivosmart 5)", [
        ("Sleep Architecture", ["TST, SE, WASO, SOL", "REM%, Deep%, Light%",
                                "N_awakenings, sleep midpoint"]),
        ("Circadian Rhythm", ["IS, IV, M10, L5, Relative Amplitude",
                               "Cosinor: amplitude, acrophase, mesor"]),
        ("Daily Summaries", ["Resting HR, nocturnal HR dip",
                              "SpO\u2082 (mean/min, T90%), stress",
                              "Daily steps, active calories"]),
    ], C_TEAL),
    (axes[1], "CGM (Dexcom G6)", [
        ("Glycemic Control", ["Mean glucose, SD, CV",
                               "5-level TIR (TBR\u2082/TBR\u2081/TIR/TAR\u2081/TAR\u2082)"]),
        ("Risk Indices", ["GMI, GRI, LBGI/HBGI",
                           "MAGE (mean amplitude of excursions)"]),
        ("Circadian Glucose", ["AGP percentile curves (p10\u2013p90)",
                                "Dawn phenomenon, nocturnal nadir",
                                "Nocturnal mean, nocturnal CV"]),
    ], C_GOLD),
]:
    ax.set_xlim(0, 10); ax.set_ylim(-0.5, 12); ax.axis("off")
    ax.text(0.2, 11.5, title, fontsize=14, fontweight="bold", color=col)
    ax.axhline(y=11.1, color=C_LGREY, linewidth=0.6, xmin=0.02, xmax=0.98)
    y = 10.5
    for section_title, bullets in items:
        ax.text(0.2, y, section_title, fontsize=12, fontweight="bold", color=C_CHAR)
        y -= 0.7
        for bullet in bullets:
            ax.text(0.6, y, "\u2022  " + bullet, fontsize=10.5, color=C_GREY)
            y -= 0.7
        y -= 0.4

fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.5, 1.3, 12.0)

# ========================== SLIDE 17: WHAT MAKES AI-READI UNIQUE ===========
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "What Makes AI-READI Unique")

points = [
    ("Dense Cross-Modal Snapshot", C_BLUE,
     "~10-day concurrent wearable + CGM + environment, anchored to same-day retinal imaging\n"
     "(4 sub-modalities, 4+ devices) + 12-lead ECG + 105 clinical labs per person."),
    ("Triple-Balanced Cohort", C_TEAL,
     "2,280 participants balanced across sex, race/ethnicity, and 4 T2DM severity groups.\n"
     "3 geographically diverse sites. Recommended train/val/test splits provided."),
    ("AI-Ready Standards", C_PURP,
     "OMOP CDM + WFDB + DICOM + Open mHealth + CDS v0.1.1. Manifests with rich metadata.\n"
     "All data de-identified, ethically documented (healthsheet.md), DOI-registered."),
    ("Metabolic Aging Lens", C_GOLD,
     "Salutogenesis framing: the pathway from T2DM back to health. Captures the full\n"
     "metabolic-to-vascular continuum with actionable continuous monitoring modalities."),
]

fig, ax = plt.subplots(figsize=(11.5, 5.2))
ax.set_xlim(0, 12); ax.set_ylim(-0.5, len(points) * 2.1); ax.axis("off")

for i, (title, col, desc) in enumerate(points):
    y = (len(points) - 1 - i) * 2.1
    # Number badge
    circle = plt.Circle((0.45, y + 0.7), 0.32, facecolor=col, alpha=0.1,
                         edgecolor=col, linewidth=1.0)
    ax.add_patch(circle)
    ax.text(0.45, y + 0.7, str(i+1), fontsize=15, fontweight="bold",
            color=col, ha="center", va="center")
    ax.text(1.1, y + 1.1, title, fontsize=14, fontweight="bold", color=col, va="center")
    ax.text(1.1, y + 0.25, desc, fontsize=10.5, color=C_GREY, va="center")

fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.5, 1.2, 12.0)

# ========================== SLIDE 18: SCALE CONTEXT ========================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Scale in Context",
              "AI-READI\u2019s role is alignment & evaluation, not from-scratch pretraining")

fig, ax = plt.subplots(figsize=(11, 4.2))
datasets = [
    ("AI-READI\n(this work)", 2280, C_RED),
    ("UK Biobank\n(imaging)", 100000, C_GOLD),
    ("GluFormer\n(CGM)", 10000000, C_TEAL),
    ("RETFound\n(retinal)", 1600000, C_PURP),
    ("ECGFounder\n(ECG)", 10000000, C_BLUE),
]
names = [d[0] for d in datasets]
vals  = [d[1] for d in datasets]
cols  = [d[2] for d in datasets]

bars = ax.bar(range(len(datasets)), vals, color=cols, width=0.45,
              edgecolor="white", linewidth=0.5, alpha=0.75)
ax.set_yscale("log"); ax.set_xticks(range(len(datasets)))
ax.set_xticklabels(names, fontsize=12)
ax.set_ylabel("Sample Size (log scale)", fontsize=12)
ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
for bar, val in zip(bars, vals):
    label = f"{val:,}" if val < 1e6 else f"{val/1e6:.0f}M"
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() * 1.5,
            label, ha="center", fontsize=11, color=C_CHAR, fontweight="bold")
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.8, 1.3, 11.5)

_bullets(sl, 0.7, 5.5, 11, 1.5, [
    "\u2022  2,280 participants \u2248 two orders of magnitude below typical pretraining corpora",
    "\u2022  Realistic role: fine-tuning pretrained models + cross-modal alignment + evaluation",
    "\u2022  Unique value: multimodal density per person, not total sample size",
], size=15, colour=GREY60)

# ========================== SLIDE 19: INFRASTRUCTURE =======================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Analysis Infrastructure (~2,700 lines Python)")

fig, ax = plt.subplots(figsize=(12, 5.0))
ax.set_xlim(0, 12); ax.set_ylim(0, 8); ax.axis("off")

layers = [
    (0.3, 6.0, 11.4, 1.5, "Layer 1: Core Data Loaders",
     "clinical.py  \u00b7  ecg.py  \u00b7  cgm.py  \u00b7  wearable.py  \u00b7  environment.py  \u00b7  retinal.py\n"
     "concepts.py  \u00b7  temporal.py  \u00b7  config.py  \u00b7  participants.py",
     C_BLUE, 1.0),
    (0.3, 4.1, 11.4, 1.5, "Layer 2: Unified Participant View",
     "participant_index.py \u2192 Parquet join of all manifests (one row per person)\n"
     "multimodal.py \u2192 get_participant(id) lazy accessor + aligned_timeseries()",
     C_TEAL, 1.0),
    (0.3, 2.2, 11.4, 1.5, "Layer 3: Feature Engineering & Analysis",
     "features.py (2,280 \u00d7 125 clinical matrix)  \u00b7  features_wearable.py (sleep, circadian, daily)\n"
     "cohort.py (ANOVA/KW, pairwise, Cohen\u2019s d, BH-FDR, covariate adjustment)",
     C_GOLD, 1.0),
    (0.3, 0.3, 11.4, 1.5, "Layer 4: ML Pipeline  [planned]",
     "PyTorch Datasets  \u00b7  Multi-modal batching  \u00b7  Augmentation  \u00b7  Models\n"
     "[Not yet implemented]",
     C_LGREY, 0.4),
]

for x, y, w, h, title, desc, col, alpha in layers:
    rect = mpatches.FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.12",
                                    facecolor=col, alpha=alpha * 0.06,
                                    edgecolor=col, linewidth=1.2)
    ax.add_patch(rect)
    ax.text(x + 0.25, y + h - 0.35, title, fontsize=13, fontweight="bold", color=col, va="top")
    ax.text(x + 0.25, y + 0.35, desc, fontsize=10, color=C_GREY, va="bottom")

# Down arrows
for yy in [6.0, 4.1, 2.2]:
    ax.annotate("", xy=(6, yy - 0.05), xytext=(6, yy + 0.05),
                arrowprops=dict(arrowstyle="->", color=C_LGREY, lw=1.2))

fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.3, 1.2, 12.5)

# ========================== SLIDE 20: RESEARCH DIRECTIONS ==================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_slide_header(sl, "Research Directions")

directions = [
    ("\u2460  Cross-Modal Aging State Vector", C_BLUE,
     "Compute full cross-modal correlation structure (CGM \u00d7 HR \u00d7 environment) during the overlap window.\n"
     "Cluster into data-driven physiological subtypes. Compare against T2DM strata."),
    ("\u2461  Non-Invasive Biomarker Prediction", C_TEAL,
     "Can continuous monitoring (wearable + CGM + env) predict expensive clinical labs?\n"
     "Minimal feature sets for HbA1c, lipids, renal function from cheap sensors."),
    ("\u2462  Cross-Modal Coupling as Aging Signal", C_GOLD,
     "Glucose\u2013HR coupling, sleep\u2013glucose interaction, environment\u2013physiology correlations.\n"
     "Does coupling strength decay with metabolic severity?"),
    ("\u2463  Multi-Agent Analysis Framework", C_PURP,
     "Domain-specialized AI agents (clinical, wearable, glucose, cardiac, retinal, environment)\n"
     "with cross-modal reasoning, hypothesis generation, and critic validation."),
]

fig, ax = plt.subplots(figsize=(12, 5.5))
ax.set_xlim(0, 12.5); ax.set_ylim(-0.5, len(directions) * 2); ax.axis("off")
for i, (title, col, desc) in enumerate(directions):
    y = (len(directions) - 1 - i) * 2
    rect = mpatches.FancyBboxPatch((0.1, y - 0.15), 12.2, 1.6,
                                    boxstyle="round,pad=0.12",
                                    facecolor=col, alpha=0.04,
                                    edgecolor=col, linewidth=1.0)
    ax.add_patch(rect)
    ax.text(0.4, y + 1.1, title, fontsize=13, fontweight="bold", color=col, va="center")
    ax.text(0.4, y + 0.25, desc, fontsize=10.5, color=C_GREY, va="center")
fig.tight_layout()
_img(sl, fig_to_buf(fig), 0.2, 1.2, 12.8)

# ========================== SLIDE 21: THANK YOU ============================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
_accent_line(sl, 4.8, 2.6, 3.6)
_tf(sl, 0.5, 2.75, 12, 0.9, "Thank You", size=44, bold=True, colour=CHARCOAL,
    align=PP_ALIGN.CENTER)
_tf(sl, 0.5, 3.7, 12, 0.5,
    "AI-READI v3.0.0  |  aireadi.org  |  DOI: 10.60775/fairhub.3",
    size=16, colour=GREY80, align=PP_ALIGN.CENTER)
_tf(sl, 0.5, 4.6, 12, 1.5,
    "2,280 participants \u00d7 9 modalities \u00d7 ~10-day overlap window\n"
    "Full loading infrastructure + feature matrix + cohort analysis\n"
    "Next: cross-modal alignment  \u00b7  multi-agent framework  \u00b7  ML pipeline",
    size=16, colour=GREY60, align=PP_ALIGN.CENTER)
_tf(sl, 0.5, 6.3, 12, 0.4,
    "Zijian (Carl) Ma  \u00b7  Stanford TWC Lab  \u00b7  carl.ma@stanford.edu",
    size=14, colour=GREYB0, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save(str(OUT))
print(f"\nDone! Saved to: {OUT}")
print(f"  {len(prs.slides)} slides, widescreen 16:9, ICML-style white theme")
